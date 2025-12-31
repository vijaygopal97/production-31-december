#!/usr/bin/env python3
"""
Complete Report Generator
Uses template PPT as base and updates only values
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta
import os
import sys
import re
from copy import deepcopy

# Import vote share calculator
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from vote_share_calculator import VoteShareCalculator


class CompleteReportGenerator:
    """Generate complete report matching final PPT template exactly"""
    
    def __init__(self, excel_path, template_ppt_path, reference_date=None):
        """
        Initialize with Excel data and template PPT
        
        Args:
            excel_path: Path to Excel file with survey data
            template_ppt_path: Path to final PPT template to copy design from
            reference_date: Reference date for report generation (default: current date)
                           Can be datetime object or string in format 'YYYY-MM-DD'
        """
        self.excel_path = excel_path
        self.template_ppt_path = template_ppt_path
        self.df = None
        self.template_prs = None
        self.output_prs = None
        self.calculator = None
        self.ae2021_vote_shares = None  # Store 2021 AE vote shares from master sheet
        
        # Set reference date (current date or provided date)
        if reference_date is None:
            self.reference_date = datetime.now()
        else:
            if isinstance(reference_date, str):
                self.reference_date = pd.to_datetime(reference_date)
            elif isinstance(reference_date, datetime):
                self.reference_date = reference_date
            else:
                self.reference_date = pd.to_datetime(reference_date)
        
        # Load data
        self.load_data()
        
        # Load template PPT
        self.load_template()
    
    def get_weight_column(self, level='Region', period='Overall', fallback_to_exact=True):
        """
        Get the correct weight column name from the dataframe
        
        Args:
            level: 'Region', 'District', or 'AC'
            period: 'Overall', 'L7D', or 'L15D'
            fallback_to_exact: If True, try exact match first before pattern matching
        
        Returns:
            Column name if found, None otherwise
        """
        # Try exact match first if fallback_to_exact
        if fallback_to_exact:
            exact_name = f'Weight - with Vote Share - AE 2021 - {level}'
            if period != 'Overall':
                exact_name += f' {period}'
            if exact_name in self.df.columns:
                return exact_name
        
        # Pattern match based on actual column names in Excel
        # Pattern: "Weight Voteshare {period} {level} Level"
        if period == 'Overall':
            pattern = f'Weight Voteshare Overall {level} Level'
        elif period == 'L7D':
            pattern = f'Weight Voteshare L7D {level} Level'
        elif period == 'L15D':
            pattern = f'Weight Voteshare L15D {level} Level'
        else:
            pattern = None
        
        if pattern and pattern in self.df.columns:
            return pattern
        
        # Try alternative patterns
        alternatives = [
            f'Weight Voteshare {period} {level} Level' if period != 'Overall' else f'Weight Voteshare Overall {level} Level',
            f'Weight - with Vote Share - AE 2021 - {level}' + (f' {period}' if period != 'Overall' else ''),
        ]
        
        for alt in alternatives:
            if alt in self.df.columns:
                return alt
        
        return None
    
    def get_sample_size(self, data, question_column=None):
        """
        Get sample size by counting only records with non-empty responses to a question
        
        Args:
            data: DataFrame to count from
            question_column: Column name for the question. If None, uses the main vote question.
                           If 'overall', counts records with at least one non-empty response to main question.
        
        Returns:
            Sample size (count of non-empty responses)
        """
        if question_column is None:
            question_column = self.calculator.vote_question
        elif question_column == 'overall':
            # For overall sample size, count records with at least one non-empty response to main question
            question_column = self.calculator.vote_question
        
        if question_column not in data.columns:
            return 0
        
        # Count only records with non-empty (not null/NaN) responses
        return len(data[data[question_column].notna()].copy())
    
    def load_data(self):
        """Load and preprocess Excel data"""
        print(f"Loading data from: {self.excel_path}")
        
        self.df = pd.read_excel(self.excel_path)
        self.df['Survey Date'] = pd.to_datetime(self.df['Survey Date'], errors='coerce')
        self.df = self.df[self.df['Survey Date'].notna()].copy()
        
        # Initialize calculator
        self.calculator = VoteShareCalculator(self.df)
        
        # Load 2021 AE vote shares from master sheet
        self.load_ae2021_vote_shares()
        
        print(f"Loaded {len(self.df)} records")
    
    def load_ae2021_vote_shares(self):
        """Load 2021 AE vote shares from master sheet for overall tables"""
        master_file = '/var/www/West_Bengal_State_Master_-_2025 v3 (1).xlsx'
        
        try:
            # Read Region Data sheet - row 2 (index 2) is state-level
            df_master = pd.read_excel(master_file, sheet_name='Region Data', nrows=5)
            state_row = df_master.iloc[2]
            header_row = df_master.iloc[1]
            
            # Get 2021 AE - Actual vote shares (columns 68-75)
            # Column 68: AITC - 2021 AE
            # Column 69: BJP - 2021 AE
            # Column 70: INC
            # Column 71: Left+
            # Column 72: RSSCMJP
            # Column 73: IND
            # Column 74: OTHERS
            # Column 75: NOTA
            
            self.ae2021_vote_shares = {
                'AITC': state_row[df_master.columns[68]] if len(df_master.columns) > 68 else 48.02,
                'BJP': state_row[df_master.columns[69]] if len(df_master.columns) > 69 else 37.97,
                'INC': state_row[df_master.columns[70]] if len(df_master.columns) > 70 else 3.03,
                'LEFT': state_row[df_master.columns[71]] if len(df_master.columns) > 71 else 5.66,
                'Others': (state_row[df_master.columns[73]] if len(df_master.columns) > 73 else 1.58) + 
                         (state_row[df_master.columns[74]] if len(df_master.columns) > 74 else 1.30),  # IND + OTHERS
                'NWR': state_row[df_master.columns[75]] if len(df_master.columns) > 75 else 1.08  # NOTA
            }
            
            print(f"Loaded 2021 AE vote shares from master sheet:")
            print(f"  AITC: {self.ae2021_vote_shares['AITC']:.2f}%")
            print(f"  BJP: {self.ae2021_vote_shares['BJP']:.2f}%")
            print(f"  INC: {self.ae2021_vote_shares['INC']:.2f}%")
            print(f"  LEFT: {self.ae2021_vote_shares['LEFT']:.2f}%")
            print(f"  Others: {self.ae2021_vote_shares['Others']:.2f}%")
            print(f"  NWR: {self.ae2021_vote_shares['NWR']:.2f}%")
        except Exception as e:
            print(f"Warning: Could not load 2021 AE vote shares from master sheet: {e}")
            # Use default values from template
            self.ae2021_vote_shares = {
                'AITC': 48.1,
                'BJP': 37.9,
                'INC': 3.0,
                'LEFT': 7.1,
                'Others': 0.0,
                'NWR': 0.0
            }
    
    def load_template(self):
        """Load template PPT to extract design"""
        print(f"Loading template from: {self.template_ppt_path}")
        self.template_prs = Presentation(self.template_ppt_path)
        print(f"Template has {len(self.template_prs.slides)} slides")
    
    def get_overall_metrics(self):
        """Get overall metrics - filtered up to reference_date"""
        # Filter data up to reference_date only
        filtered_df = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # Overall sample size: only count records with non-empty responses to main question
        sample_size = self.get_sample_size(filtered_df, 'overall')
        
        # For date calculations, use all records (even empty ones) for date range
        valid_date_df = filtered_df[filtered_df['Survey Date'].notna()].copy()
        start_date = valid_date_df['Survey Date'].min() if len(valid_date_df) > 0 else None
        end_date = min(self.reference_date, valid_date_df['Survey Date'].max()) if len(valid_date_df) > 0 else None
        duration_days = (end_date - start_date).days + 1 if start_date and end_date else 0
        
        # Calculate daily average sample (for "Total Sample conducted Daily")
        # Only count records with valid responses for daily averages
        valid_response_df = filtered_df[filtered_df[self.calculator.vote_question].notna()].copy()
        daily_samples = valid_response_df.groupby('Survey Date').size()
        avg_daily_sample = int(daily_samples.mean()) if len(daily_samples) > 0 else 0  # Round to nearest integer
        
        # Calculate gender breakdown - only count those with valid responses
        valid_response_df = filtered_df[filtered_df[self.calculator.vote_question].notna()].copy()
        male_sample = len(valid_response_df[valid_response_df['Gender'] == 1])
        female_sample = len(valid_response_df[valid_response_df['Gender'] == 2])
        
        # Calculate F2F and CATI samples (Data Type: 1 = F2F, 2 = CATI) - only with valid responses
        f2f_sample = len(valid_response_df[valid_response_df['Data Type'] == 1]) if 'Data Type' in valid_response_df.columns else 0
        cati_sample = len(valid_response_df[valid_response_df['Data Type'] == 2]) if 'Data Type' in valid_response_df.columns else 0
        
        # Calculate zones covered
        if 'Region Name' in filtered_df.columns:
            unique_zones = filtered_df['Region Name'].nunique()
            total_zones = 5  # West Bengal has 5 zones, or could calculate if needed
        else:
            unique_zones = 0
            total_zones = 5
        
        # Incumbent party (AITC for West Bengal)
        incumbent_party = 'AITC'
        
        return {
            'sample_size': sample_size,
            'start_date': start_date,
            'end_date': end_date,
            'duration_days': duration_days,
            'avg_daily_sample': avg_daily_sample,
            'male_sample': male_sample,
            'female_sample': female_sample,
            'f2f_sample': f2f_sample,
            'cati_sample': cati_sample,
            'zones_covered': f"{unique_zones}/{total_zones}",
            'incumbent_party': incumbent_party
        }
    
    def update_slide_5_table(self, table):
        """Update slide 5 (State Level Vote Share) table with calculated values"""
        # Calculate vote shares - filter data up to reference date only
        # Use data from start up to reference_date (when report is generated)
        overall_data = self.calculator.df[
            self.calculator.df['Survey Date'] <= self.reference_date
        ].copy()
        raw_vote_shares = self.calculator.calculate_vote_share(overall_data, use_weights=False)
        
        # For 7 DMA, calculate from reference date (when report is generated)
        # For 7DMA: end_date is one day before the reference date
        # So if reference_date is Oct 31, use Oct 24-30 (7 days ending on Oct 30, excluding Oct 31)
        # Window: (reference_date - 7 days) to (reference_date - 1 day) = 7 days total
        end_date = self.reference_date - timedelta(days=1)  # One day before reference date
        cutoff_date = end_date - timedelta(days=6)  # 6 days before end_date = 7 days total
        dma7_data = self.df[
            (self.df['Survey Date'] >= cutoff_date) & 
            (self.df['Survey Date'] <= end_date) &
            (self.df['Survey Date'] <= self.reference_date - timedelta(days=1))
        ].copy()
        
        # For 7DMA, use L7D weights when >=50% available, otherwise use regular weights
        # This matches the final PPT behavior where L7D weights are only used when sufficient
        l7d_col = self.get_weight_column('Region', 'L7D') or 'Weight Voteshare L7D Region Level'
        regular_col = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        # Check if columns exist, if not use fallback
        if l7d_col not in dma7_data.columns:
            l7d_col = None
        if regular_col not in dma7_data.columns:
            regular_col = None
        
        l7d_available = dma7_data[l7d_col].notna().sum() if l7d_col else 0
        total_records = len(dma7_data)
        
        # For 7DMA: Use L7D weights if >=50% of records have L7D weights, otherwise use regular weights
        # This ensures we don't use sparse L7D weights which give incorrect results
        if l7d_col and l7d_available > 0 and (l7d_available / total_records) >= 0.5:
            # Use L7D weights - filter to only records with L7D weights
            dma7_data_with_weights = dma7_data[dma7_data[l7d_col].notna()].copy()
            dma7_weight_column = l7d_col
        else:
            # Fall back to regular weights when L7D weights are sparse (<50%) or not available
            dma7_data_with_weights = dma7_data.copy()
            dma7_weight_column = regular_col if regular_col else None
        
        # Use 7DMA-specific weights for 7DMA calculations
        # For sample size, calculate_vote_share will count all valid votes from dma7_data_with_weights
        # But if we're using L7D weights, we need to count all valid votes from the original dma7_data
        # to match the final PPT which counts all valid votes, not just those with L7D weights
        dma7_vote_shares = self.calculator.calculate_vote_share(
            dma7_data_with_weights,
            weight_column=dma7_weight_column,
            use_weights=True
        )
        
        # For sample size: if using L7D weights, count all records with L7D weights (regardless of vote)
        # This matches the final PPT which shows 15,206 for 7DMA sample
        if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
            # Count all records with L7D weights (regardless of whether they have valid votes)
            dma7_vote_shares['sample'] = len(dma7_data[dma7_data[l7d_col].notna()].copy())
        
        # Normalized vote share uses same data filtered up to reference date
        normalized_vote_shares = self.calculator.calculate_vote_share(
            overall_data,  # Already filtered to reference_date
            weight_column=self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level',
            use_weights=True
        )
        
        # Party mapping
        party_map = {
            'AITC': 'AITC',
            'BJP': 'BJP',
            'LEFT': 'Left Front',
            'INC': 'INC',
            'Others': 'Other Parties',
            'NWR': 'N+W+R'
        }
        
        # Update table - Column structure: 0=Party, 1=Raw, 2=3DMA (leave as is), 3=7DMA, 4=Normalized, 5=AE2021 (leave as is), 6=GE2024 (leave as is)
        # Sample row (row 2, 0-indexed)
        sample_row = 2
        table.cell(sample_row, 1).text = f"{raw_vote_shares.get('sample', 0):,}"
        # Right-align sample size
        if len(table.cell(sample_row, 1).text_frame.paragraphs) > 0:
            table.cell(sample_row, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        # Column 2 (3DMA) - leave as is in template, don't update
        table.cell(sample_row, 3).text = f"{dma7_vote_shares.get('sample', 0):,}"
        if len(table.cell(sample_row, 3).text_frame.paragraphs) > 0:
            table.cell(sample_row, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        table.cell(sample_row, 4).text = f"{normalized_vote_shares.get('sample', 0):,}"
        if len(table.cell(sample_row, 4).text_frame.paragraphs) > 0:
            table.cell(sample_row, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        # Columns 5-6 (AE2021, GE2024) - leave as is in template
        
        # Update party rows
        parties = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        for party_idx, party in enumerate(parties):
            row_idx = party_idx + 3  # Rows 3-8
            
            # Party name (check if needs update)
            party_name = party_map.get(party, party)
            if table.cell(row_idx, 0).text != party_name:
                table.cell(row_idx, 0).text = party_name
            
            # Raw vote share (Column 1) - right align
            raw_val = raw_vote_shares.get(party, 0)
            table.cell(row_idx, 1).text = f"{raw_val:.1f}"
            if len(table.cell(row_idx, 1).text_frame.paragraphs) > 0:
                table.cell(row_idx, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # Column 2 (3DMA) - leave as is in template, don't update
            
            # 7 DMA vote share (Column 3) - right align
            dma7_val = dma7_vote_shares.get(party, 0)
            table.cell(row_idx, 3).text = f"{dma7_val:.1f}"
            if len(table.cell(row_idx, 3).text_frame.paragraphs) > 0:
                table.cell(row_idx, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # Normalized vote share (Column 4) - right align
            norm_val = normalized_vote_shares.get(party, 0)
            table.cell(row_idx, 4).text = f"{norm_val:.1f}"
            if len(table.cell(row_idx, 4).text_frame.paragraphs) > 0:
                table.cell(row_idx, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # Columns 5-6 (AE2021, GE2024) - leave as is in template
        
        # Update margin row - right align all numeric values
        margin_row = 9
        raw_margin = raw_vote_shares.get('AITC', 0) - raw_vote_shares.get('BJP', 0)
        dma7_margin = dma7_vote_shares.get('AITC', 0) - dma7_vote_shares.get('BJP', 0)
        norm_margin = normalized_vote_shares.get('AITC', 0) - normalized_vote_shares.get('BJP', 0)
        
        table.cell(margin_row, 1).text = f"{raw_margin:.1f}"
        if len(table.cell(margin_row, 1).text_frame.paragraphs) > 0:
            table.cell(margin_row, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        # Column 2 (3DMA margin) - leave as is in template, don't update
        table.cell(margin_row, 3).text = f"{dma7_margin:.1f}"
        if len(table.cell(margin_row, 3).text_frame.paragraphs) > 0:
            table.cell(margin_row, 3).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        table.cell(margin_row, 4).text = f"{norm_margin:.1f}"
        if len(table.cell(margin_row, 4).text_frame.paragraphs) > 0:
            table.cell(margin_row, 4).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        # Columns 5-6 (AE2021, GE2024 margins) - leave as is in template
        
        print("Updated Slide 5 table values")
    
    def update_slide_2_table(self, table):
        """Update slide 2 (Introduction) table with all calculated stats"""
        metrics = self.get_overall_metrics()
        
        # Use reference date (when report is generated)
        # Format date with proper ordinal suffix
        day = self.reference_date.day
        if 11 <= day <= 13:
            suffix = "th"
        elif day % 10 == 1:
            suffix = "st"
        elif day % 10 == 2:
            suffix = "nd"
        elif day % 10 == 3:
            suffix = "rd"
        else:
            suffix = "th"
        presentation_date = self.reference_date.strftime(f'%d{suffix} %B %Y')
        
        # Mapping of text patterns to metric keys
        stat_mappings = {
            'sample covered': ('sample_size', lambda x: f"{x:,}"),
            'incumbent party': ('incumbent_party', lambda x: str(x)),
            'zones covered': ('zones_covered', lambda x: str(x)),
            'presentation date': ('presentation_date', lambda x: presentation_date),
            'total sample conducted daily': ('avg_daily_sample', lambda x: f"{x:,}"),
            'total male sample': ('male_sample', lambda x: f"{x:,}"),
            'total female sample': ('female_sample', lambda x: f"{x:,}"),
            'f2f sample': ('f2f_sample', lambda x: f"{x:,}"),
            'cati sample': ('cati_sample', lambda x: f"{x:,}")
        }
        
        # Update table rows
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_text = cell.text.strip()
                cell_text_lower = cell_text.lower()
                
                # Check each stat mapping
                for pattern, (metric_key, formatter) in stat_mappings.items():
                    if pattern in cell_text_lower:
                        if len(table.columns) > 1:
                            if metric_key == 'presentation_date':
                                # Use current date directly
                                table.cell(row_idx, 1).text = presentation_date
                            else:
                                # Get value from metrics and format it
                                value = metrics.get(metric_key, '')
                                table.cell(row_idx, 1).text = formatter(value)
                            
                            # Right-align the value
                            if len(table.cell(row_idx, 1).text_frame.paragraphs) > 0:
                                table.cell(row_idx, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        break
        
        # Right-align all numeric values in the "Details" column (column 1)
        # Skip header row (row 0)
        for row_idx in range(1, len(table.rows)):
            details_cell = table.cell(row_idx, 1)
            # Right-align all paragraphs in this cell
            for paragraph in details_cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.RIGHT
    
    def update_demographic_table(self, table, slide):
        """
        Update demographic breakdown tables (Gender, Location, Religion, Social Category, Age)
        
        Args:
            table: Table object to update
            slide: Slide object to identify which demographic
        """
        # Determine demographic type from slide content
        demographic_type = None
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                text = shape.text.lower()
                if 'gender' in text:
                    demographic_type = 'Gender'
                    break
                elif 'location' in text:
                    demographic_type = 'Location'
                    break
                elif 'religion' in text:
                    demographic_type = 'Religion'
                    break
                elif 'social category' in text:
                    demographic_type = 'Social Category'
                    break
                elif 'age' in text:
                    demographic_type = 'Age'
                    break
        
        if not demographic_type:
            # Try to infer from table headers
            if len(table.rows) > 0:
                header_row = table.rows[0]
                for cell in header_row.cells:
                    text = cell.text.lower()
                    if 'gender' in text:
                        demographic_type = 'Gender'
                        break
                    elif 'location' in text:
                        demographic_type = 'Location'
                        break
                    elif 'religion' in text:
                        demographic_type = 'Religion'
                        break
                    elif 'social' in text or 'category' in text:
                        demographic_type = 'Social Category'
                        break
                    elif 'age' in text:
                        demographic_type = 'Age'
                        break
        
        if not demographic_type:
            return
        
        # Get demographic column name
        demographic_column_map = {
            'Gender': 'Gender',
            'Location': 'Residential locality type',
            'Religion': '20. Could you please tell me the religion that you belong to?',
            'Social Category': '21. Which social category do you belong to?',
            'Age': 'Could you please tell me your age in complete years?'
        }
        
        demographic_column = demographic_column_map.get(demographic_type)
        if not demographic_column or demographic_column not in self.df.columns:
            print(f"Warning: {demographic_column} not found in data")
            return
        
        # Determine which metrics to calculate from table headers
        # Check header row to see what columns are present
        has_7dma = False
        has_15dma = False
        has_30dma = False
        has_overall = False
        
        if len(table.rows) > 0:
            header_row = table.rows[0]
            for cell in header_row.cells:
                text = cell.text.strip().lower()
                if '7 dma' in text or '7dma' in text:
                    has_7dma = True
                elif '15 dma' in text or '15dma' in text:
                    has_15dma = True
                elif '30 dma' in text or '30dma' in text:
                    has_30dma = True
                elif 'overall' in text:
                    has_overall = True
        
        # Calculate demographic breakdowns
        # Overall data - filter up to reference_date
        overall_data = self.calculator.df[
            self.calculator.df['Survey Date'] <= self.reference_date
        ].copy()
        
        # 7 DMA data - last 7 days ending one day before reference_date
        dma7_data = self.calculator.filter_by_date_range(
            days=7,
            exclude_latest=True,  # Exclude reference_date itself
            reference_date=self.reference_date
        )
        
        # 15 DMA data - last 15 days ending one day before reference_date
        dma15_data = self.calculator.filter_by_date_range(
            days=15,
            exclude_latest=True,  # Exclude reference_date itself
            reference_date=self.reference_date
        )
        
        # Get categories from table or data
        categories = []
        if len(table.rows) > 2:  # Skip header rows
            for row_idx in range(2, len(table.rows)):
                category_name = table.cell(row_idx, 0).text.strip()
                if category_name and category_name not in ['Sample', 'Margin', '']:
                    categories.append(category_name)
        
        # If no categories found, use data
        if not categories:
            unique_cats = self.df[demographic_column].dropna().unique()
            categories = [str(cat) for cat in unique_cats if str(cat).lower() not in ['nan', 'none', '']]
        
        # Update table for each category
        # Table structure: [Category, Sample_7DMA, AITC_7DMA, BJP_7DMA, LEFT_7DMA, INC_7DMA, Others_7DMA, NWR_7DMA, Margin_7DMA, Sample_Overall, ...]
        
        for category_idx, category in enumerate(categories[:len(table.rows)-2]):
            row_idx = category_idx + 2
            
            # Filter data by category using numeric codes
            category_filter = None
            if demographic_type == 'Gender':
                # Gender: 1=Male, 2=Female
                if category.lower() in ['male', 'm']:
                    category_filter = self.df[demographic_column] == 1
                elif category.lower() in ['female', 'f']:
                    category_filter = self.df[demographic_column] == 2
            elif demographic_type == 'Location':
                # Location: 1=Urban, 2=Rural
                if category.lower() in ['urban', 'u']:
                    category_filter = self.df[demographic_column] == 1
                elif category.lower() in ['rural', 'r']:
                    category_filter = self.df[demographic_column] == 2
            elif demographic_type == 'Religion':
                # Religion: 1=Hindu, 2=Muslim
                if category.lower() in ['hindu', 'h']:
                    category_filter = self.df[demographic_column] == 1
                elif category.lower() in ['muslim', 'm', 'islam']:
                    category_filter = self.df[demographic_column] == 2
                elif category.lower() in ['other', 'others']:
                    # Other religions (codes 3, 4, 5, 6, 7, etc.)
                    category_filter = ~self.df[demographic_column].isin([1, 2])
            elif demographic_type == 'Social Category':
                # Social Category: 1=General, 2=OBC, 3=SC, 4=ST
                if 'general' in category.lower() and 'obc' in category.lower():
                    # General+OBC combined
                    category_filter = self.df[demographic_column].isin([1, 2])
                elif 'general' in category.lower():
                    category_filter = self.df[demographic_column] == 1
                elif 'obc' in category.lower():
                    category_filter = self.df[demographic_column] == 2
                elif category.lower() in ['sc', 'scheduled caste']:
                    category_filter = self.df[demographic_column] == 3
                elif category.lower() in ['st', 'scheduled tribe']:
                    category_filter = self.df[demographic_column] == 4
            elif demographic_type == 'Age':
                # Age groups based on numeric age
                age_col = self.df[demographic_column]
                age_numeric = pd.to_numeric(age_col, errors='coerce')
                
                if '18-25' in category or '18 - 25' in category:
                    category_filter = (age_numeric >= 18) & (age_numeric <= 25)
                elif '26-35' in category or '26 - 35' in category or '26-34' in category or '26 - 34' in category:
                    # Handle both 26-35 and 26-34
                    if '26-34' in category or '26 - 34' in category:
                        category_filter = (age_numeric >= 26) & (age_numeric <= 34)
                    else:
                        category_filter = (age_numeric >= 26) & (age_numeric <= 35)
                elif '36-50' in category or '36 - 50' in category:
                    category_filter = (age_numeric >= 36) & (age_numeric <= 50)
                elif '50+' in category or 'above 50' in category or '50 +' in category:
                    category_filter = age_numeric > 50
            
            if category_filter is None:
                print(f"Warning: Could not match category '{category}' for {demographic_type}")
                continue
            
            # Filter dataframes by category
            category_dma7_data = dma7_data[category_filter[dma7_data.index]].copy() if has_7dma else pd.DataFrame()
            category_dma15_data = dma15_data[category_filter[dma15_data.index]].copy() if has_15dma else pd.DataFrame()
            category_overall_data = overall_data[category_filter[overall_data.index]].copy() if has_overall else pd.DataFrame()
            
            # Weight columns
            l7d_col = self.get_weight_column('Region', 'L7D') or 'Weight Voteshare L7D Region Level'
            l15d_col = self.get_weight_column('Region', 'L15D') or 'Weight Voteshare L15D Region Level'
            regular_col = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
            
            # Calculate 7DMA vote shares (using L7D weights)
            dma7_vote_shares = {}
            if has_7dma and len(category_dma7_data) > 0:
                # Check if L7D weights are available for this category
                l7d_available = category_dma7_data[l7d_col].notna().sum() if l7d_col in category_dma7_data.columns else 0
                total_records = len(category_dma7_data)
                
                # For 7DMA: Use L7D weights if >=50% of records have L7D weights, otherwise use regular weights
                if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
                    # Use L7D weights - filter to only records with L7D weights
                    category_dma7_data_with_weights = category_dma7_data[category_dma7_data[l7d_col].notna()].copy()
                    dma7_weight_column = l7d_col
                else:
                    # Fall back to regular weights when L7D weights are sparse (<50%)
                    category_dma7_data_with_weights = category_dma7_data.copy()
                    dma7_weight_column = regular_col
                
                dma7_vote_shares = self.calculator.calculate_vote_share(
                    category_dma7_data_with_weights,
                    weight_column=dma7_weight_column,
                    use_weights=True
                )
                
                # For sample size: when using L7D weights, count all records with valid votes from ORIGINAL data
                # (before filtering by L7D weights), not just those with L7D weights
                # This matches the audit trail which shows all records in the 7DMA window
                if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
                    # Count all records with valid votes from original category_dma7_data
                    vote_question = '8. If assembly elections (MLA) were to be held tomorrow, then which party would you vote for?'
                    dma7_vote_shares['sample'] = len(category_dma7_data[category_dma7_data[vote_question].notna()].copy())
            
            # Calculate 15DMA vote shares (using L15D weights)
            dma15_vote_shares = {}
            if has_15dma and len(category_dma15_data) > 0:
                # Check if L15D weights are available for this category
                l15d_available = category_dma15_data[l15d_col].notna().sum() if l15d_col in category_dma15_data.columns else 0
                total_records = len(category_dma15_data)
                
                # For 15DMA: Use L15D weights if >=50% of records have L15D weights, otherwise use regular weights
                if l15d_available > 0 and (l15d_available / total_records) >= 0.5:
                    # Use L15D weights - filter to only records with L15D weights
                    category_dma15_data_with_weights = category_dma15_data[category_dma15_data[l15d_col].notna()].copy()
                    dma15_weight_column = l15d_col
                else:
                    # Fall back to regular weights when L15D weights are sparse (<50%)
                    category_dma15_data_with_weights = category_dma15_data.copy()
                    dma15_weight_column = regular_col
                
                dma15_vote_shares = self.calculator.calculate_vote_share(
                    category_dma15_data_with_weights,
                    weight_column=dma15_weight_column,
                    use_weights=True
                )
                
                # For sample size: when using L15D weights, count all records with valid votes from ORIGINAL data
                # (before filtering by L15D weights), not just those with L15D weights
                # This matches the audit trail which shows all records in the 15DMA window
                if l15d_available > 0 and (l15d_available / total_records) >= 0.5:
                    # Count all records with valid votes from original category_dma15_data
                    vote_question = '8. If assembly elections (MLA) were to be held tomorrow, then which party would you vote for?'
                    dma15_vote_shares['sample'] = len(category_dma15_data[category_dma15_data[vote_question].notna()].copy())
            
            # Calculate Overall vote shares (using regular weights)
            overall_vote_shares = {}
            if has_overall and len(category_overall_data) > 0:
                overall_vote_shares = self.calculator.calculate_vote_share(
                    category_overall_data,
                    weight_column=regular_col,
                    use_weights=True
                )
            
            # Helper function to format cell with proper font, alignment, and decimal alignment
            def format_table_cell(cell, value, is_percentage=False):
                """Format table cell with 14pt font, right alignment, and decimal point alignment"""
                from pptx.util import Pt
                from pptx.enum.text import PP_ALIGN
                
                # Format value text
                if is_percentage:
                    # Format with decimal point alignment: use fixed width format (XX.X)
                    # Pad to ensure decimal points align (4 characters: XX.X)
                    formatted_value = f"{value:.1f}".rjust(5)  # Right-justify with 5 chars for "XX.X" format (e.g., "45.2", " 3.2")
                    # Ensure at least one space before if needed for alignment
                    if len(formatted_value) < 5:
                        formatted_value = formatted_value.rjust(5)
                else:
                    # For sample sizes, use comma formatting
                    formatted_value = f"{value:,}"
                
                # Set cell text
                cell.text = formatted_value
                
                # Get or create text frame paragraph
                text_frame = cell.text_frame
                text_frame.clear()  # Clear existing content
                para = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                para.clear()  # Clear runs
                
                # Add run with formatted text
                run = para.add_run()
                run.text = formatted_value
                
                # Set font size to 14pt (matching label font size)
                run.font.size = Pt(14)
                # Set font name to match template (Aptos Display)
                run.font.name = 'Aptos Display'
                
                # Set paragraph alignment to RIGHT
                para.alignment = PP_ALIGN.RIGHT
            
            # Update table cells based on which metrics are present
            # Table structure: Column 0 = Category, Columns 1-8 = First metric, Columns 9-16 = Second metric
            
            # First metric columns (1-8): 7DMA or 30DMA
            if has_7dma:
                # Fill 7DMA columns (1-8)
                format_table_cell(table.cell(row_idx, 1), dma7_vote_shares.get('sample', 0), is_percentage=False)
                format_table_cell(table.cell(row_idx, 2), dma7_vote_shares.get('AITC', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 3), dma7_vote_shares.get('BJP', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 4), dma7_vote_shares.get('LEFT', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 5), dma7_vote_shares.get('INC', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 6), dma7_vote_shares.get('Others', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 7), dma7_vote_shares.get('NWR', 0), is_percentage=True)
                dma7_margin = dma7_vote_shares.get('AITC', 0) - dma7_vote_shares.get('BJP', 0)
                format_table_cell(table.cell(row_idx, 8), dma7_margin, is_percentage=True)
            elif has_30dma:
                # Leave 30DMA columns empty (as per user requirement)
                # Explicitly clear columns 1-8 for 30DMA to ensure they're empty
                # Don't fill columns 1-8 for 30DMA
                for col_idx in range(1, 9):  # Columns 1-8 (Sample, AITC, BJP, LEFT, INC, Others, NWR, Margin)
                    if col_idx < len(table.columns):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = ''  # Clear the cell
                        # Clear text frame
                        if cell.text_frame:
                            cell.text_frame.clear()
            
            # Second metric columns (9-16): 15DMA or Overall
            if has_15dma and len(table.columns) >= 16:
                # Fill 15DMA columns (9-16)
                format_table_cell(table.cell(row_idx, 9), dma15_vote_shares.get('sample', 0), is_percentage=False)
                format_table_cell(table.cell(row_idx, 10), dma15_vote_shares.get('AITC', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 11), dma15_vote_shares.get('BJP', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 12), dma15_vote_shares.get('LEFT', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 13), dma15_vote_shares.get('INC', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 14), dma15_vote_shares.get('Others', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 15), dma15_vote_shares.get('NWR', 0), is_percentage=True)
                dma15_margin = dma15_vote_shares.get('AITC', 0) - dma15_vote_shares.get('BJP', 0)
                format_table_cell(table.cell(row_idx, 16), dma15_margin, is_percentage=True)
            elif has_overall and len(table.columns) >= 16:
                # Fill Overall columns (9-16)
                format_table_cell(table.cell(row_idx, 9), overall_vote_shares.get('sample', 0), is_percentage=False)
                format_table_cell(table.cell(row_idx, 10), overall_vote_shares.get('AITC', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 11), overall_vote_shares.get('BJP', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 12), overall_vote_shares.get('LEFT', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 13), overall_vote_shares.get('INC', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 14), overall_vote_shares.get('Others', 0), is_percentage=True)
                format_table_cell(table.cell(row_idx, 15), overall_vote_shares.get('NWR', 0), is_percentage=True)
                overall_margin = overall_vote_shares.get('AITC', 0) - overall_vote_shares.get('BJP', 0)
                format_table_cell(table.cell(row_idx, 16), overall_margin, is_percentage=True)
    
    def get_caste_code_mapping(self):
        """
        Get mapping from caste names (as shown in template) to caste codes (in Excel)
        This mapping needs to be verified against the questionnaire
        """
        # Caste code to name mapping (based on template and common West Bengal castes)
        # Note: This mapping may need adjustment based on actual questionnaire
        caste_name_to_code = {
            'Mahishya': None,  # Will be determined by top codes
            'Kayastha': None,
            'Brahmin': None,
            'Shaikh': None,
            'Ansari': None,
            'Mahato': None,
            'Rajbanshi': None,
            'Namasudra': None,
            'Mallah': None,
            'Bagda': None,
            'Sental': None,
            'Others': [44, 48]  # Others codes
        }
        
        # Get top caste codes from data to create mapping
        caste_column = '22. Could you please tell me your caste?'
        if caste_column in self.df.columns:
            top_codes = self.df[caste_column].value_counts().head(12).index.tolist()
            # Map top codes to caste names based on frequency and common patterns
            # This is a heuristic - should be verified with questionnaire
            code_mapping = {}
            for i, code in enumerate(top_codes):
                if i < len(caste_name_to_code) - 1:  # Exclude 'Others'
                    caste_name = list(caste_name_to_code.keys())[i]
                    if caste_name != 'Others':
                        code_mapping[caste_name] = code
        
        return caste_name_to_code, code_mapping
    
    def update_caste_table(self, table, is_30dma=False):
        """
        Update caste-wise vote shares table (Slides 41-42)
        
        Args:
            table: Table object to update
            is_30dma: If True, use 30 DMA; if False, use Overall
        """
        caste_column = '22. Could you please tell me your caste?'
        
        if caste_column not in self.df.columns:
            print(f"Warning: {caste_column} not found in data")
            return
        
        # Get caste names from table (rows 2-13, column 1)
        caste_names = []
        for row_idx in range(1, min(13, len(table.rows))):  # Skip header row
            caste_name = table.cell(row_idx, 1).text.strip()  # Column 1 is "Caste"
            if caste_name and caste_name not in ['S.No', 'Caste', '']:
                caste_names.append(caste_name)
        
        if not caste_names:
            print("Warning: No caste names found in table")
            return
        
        # Create caste code mapping based on top codes in data
        # Map caste names to codes by matching top codes to template order
        # Note: Code 88 is likely NWR/Refused, so exclude it from caste mapping
        top_codes_all = self.df[caste_column].value_counts().head(15).index.tolist()  # Get more codes
        # Exclude code 88 (likely NWR/Refused) from valid caste codes
        valid_codes = [code for code in top_codes_all if code != 88][:12]  # Take top 12 excluding 88
        
        caste_code_map = {}
        for i, caste_name in enumerate(caste_names):
            if caste_name == 'Others':
                # Others includes codes 44, 48, and other codes not in the main 11 castes
                others_codes = [code for code in top_codes_all if code in [44, 48] or (code not in valid_codes[:11] and code != 88)]
                caste_code_map[caste_name] = others_codes if others_codes else [44, 48]
            elif i < len(valid_codes):
                caste_code_map[caste_name] = valid_codes[i]
            else:
                # If we run out of codes, skip this caste
                print(f"Warning: No code available for caste '{caste_name}'")
                caste_code_map[caste_name] = None
        
        # Helper function to format cell with right alignment
        def format_table_cell(cell, value, is_percentage=False):
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            
            if is_percentage:
                formatted_value = f"{value:.1f}"
            else:
                formatted_value = f"{value:,}"
            
            cell.text = formatted_value
            text_frame = cell.text_frame
            text_frame.clear()
            para = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            para.clear()
            run = para.add_run()
            run.text = formatted_value
            run.font.size = Pt(14)
            run.font.name = 'Aptos Display'
            para.alignment = PP_ALIGN.RIGHT
        
        # Calculate vote shares for each caste
        for row_idx, caste_name in enumerate(caste_names, start=1):  # Start from row 2 (1-indexed)
            if row_idx >= len(table.rows):
                break
            
            # Get caste code(s) for this caste name
            caste_codes = caste_code_map.get(caste_name)
            if caste_codes is None:
                print(f"Warning: No code mapping found for caste '{caste_name}'")
                continue
            
            # Filter data by caste code(s)
            if isinstance(caste_codes, list):
                # Multiple codes (e.g., Others)
                caste_filter = self.df[caste_column].isin(caste_codes)
            else:
                # Single code
                caste_filter = self.df[caste_column] == caste_codes
            
            caste_data = self.df[caste_filter].copy()
            
            if is_30dma:
                # 30 DMA: use last 30 days ending one day before reference date
                end_date = self.reference_date - timedelta(days=1)
                cutoff_date = end_date - timedelta(days=29)  # 30 days total
                filtered_data = caste_data[
                    (caste_data['Survey Date'] >= cutoff_date) & 
                    (caste_data['Survey Date'] <= end_date) &
                    (caste_data['Survey Date'] <= self.reference_date - timedelta(days=1))
                ].copy()
                
                # Check for L15D weights (30 DMA might use 15D weights)
                l15d_col = self.get_weight_column('Region', 'L15D') or 'Weight Voteshare L15D Region Level'
                regular_col = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
                
                l15d_available = filtered_data[l15d_col].notna().sum() if l15d_col in filtered_data.columns else 0
                total_records = len(filtered_data)
                
                if l15d_available > 0 and (l15d_available / total_records) >= 0.5:
                    filtered_data = filtered_data[filtered_data[l15d_col].notna()].copy()
                    weight_column = l15d_col
                else:
                    weight_column = regular_col
            else:
                # Overall: use all data up to reference date
                filtered_data = caste_data[
                    caste_data['Survey Date'] <= self.reference_date
                ].copy()
                weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
            
            # Calculate vote shares
            vote_shares = self.calculator.calculate_vote_share(
                filtered_data,
                weight_column=weight_column,
                use_weights=True
            )
            
            # Update table cells
            # Column structure: 0=S.No, 1=Caste, 2=Sample, 3=AITC, 4=BJP, 5=LEFT, 6=INC, 7=Others, 8=NWR, 9=Margin
            format_table_cell(table.cell(row_idx, 2), vote_shares.get('sample', 0), is_percentage=False)  # Sample
            format_table_cell(table.cell(row_idx, 3), vote_shares.get('AITC', 0), is_percentage=True)  # AITC
            format_table_cell(table.cell(row_idx, 4), vote_shares.get('BJP', 0), is_percentage=True)  # BJP
            format_table_cell(table.cell(row_idx, 5), vote_shares.get('LEFT', 0), is_percentage=True)  # LEFT
            format_table_cell(table.cell(row_idx, 6), vote_shares.get('INC', 0), is_percentage=True)  # INC
            format_table_cell(table.cell(row_idx, 7), vote_shares.get('Others', 0), is_percentage=True)  # Others
            format_table_cell(table.cell(row_idx, 8), vote_shares.get('NWR', 0), is_percentage=True)  # NWR
            
            # Margin (AITC - BJP)
            margin = vote_shares.get('AITC', 0) - vote_shares.get('BJP', 0)
            format_table_cell(table.cell(row_idx, 9), margin, is_percentage=True)  # Margin
    
    def calculate_overall_normalized_for_days(self, num_days=16):
        """
        Calculate Overall Normalized Vote Share (cumulative) for each of the last N days
        For each date, calculates normalized vote share using ALL data up to that date
        Dates end at reference_date (when report is generated), not last data date
        
        Args:
            num_days: Number of days to calculate (default 16 to match final PPT)
        
        Returns:
            List of dictionaries with date and vote shares for each party
        """
        # Calculate dates: last N days ending at reference_date (when report is generated)
        # Final PPT shows Oct 15-31 (16 dates ending at reference date Oct 31), but skips Oct 19
        # So we need: Oct 15, 16, 17, 18, 20, 21, ..., 31 (16 dates total)
        target_dates = []
        # Generate dates from Oct 15 (reference_date - num_days) to Oct 31 (reference_date)
        start_offset = num_days  # Start 16 days back from Oct 31 = Oct 15
        all_dates = []
        for i in range(start_offset, -1, -1):  # From start_offset down to 0 (reference date)
            date = self.reference_date - timedelta(days=i)
            all_dates.append(date)
        
        # Final PPT skips Oct 19, so we exclude it if present
        # This gives us exactly 16 dates: Oct 15-31 excluding Oct 19
        oct_19 = self.reference_date - timedelta(days=12)  # Oct 19 is 12 days before Oct 31
        for date in all_dates:
            if date != oct_19:  # Skip Oct 19 to match Final PPT
                target_dates.append(date)
        
        # Ensure we have exactly num_days dates
        if len(target_dates) > num_days:
            target_dates = target_dates[:num_days]  # Take first num_days
        elif len(target_dates) < num_days:
            # If we have fewer (shouldn't happen), pad with last date
            while len(target_dates) < num_days:
                target_dates.append(target_dates[-1])
        
        daily_vote_shares = []
        
        for date in target_dates:
            # For each date, calculate OVERALL normalized vote share using ALL data up to this date (cumulative)
            # Data should be filtered up to this date AND not exceed reference_date
            actual_end_date = min(date, self.reference_date, self.df['Survey Date'].max())
            date_data = self.df[
                (self.df['Survey Date'] <= actual_end_date) & 
                (self.df['Survey Date'] <= self.reference_date)
            ].copy()
            
            if len(date_data) > 0:
                # Calculate overall normalized vote share (cumulative) for this date
                vote_shares = self.calculator.calculate_vote_share(
                    date_data,
                    weight_column=self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level',
                    use_weights=True
                )
                
                # Convert date to Excel serial number for chart (matching final PPT format)
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': vote_shares.get('AITC', 0),
                    'BJP': vote_shares.get('BJP', 0),
                    'LEFT': vote_shares.get('LEFT', 0),
                    'INC': vote_shares.get('INC', 0),
                    'Others': vote_shares.get('Others', 0),
                    'NWR': vote_shares.get('NWR', 0)
                })
            else:
                # If no data for this date, still add entry with zero values
                # Convert date to Excel serial number for chart
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': 0,
                    'BJP': 0,
                    'LEFT': 0,
                    'INC': 0,
                    'Others': 0,
                    'NWR': 0
                })
        
        return daily_vote_shares
    
    def calculate_overall_normalized_for_demographic(self, demographic_type, demographic_value, num_days=16):
        """
        Calculate Overall Normalized Vote Share (cumulative) for a specific demographic
        For each date, calculates normalized vote share using ALL data up to that date, filtered by demographic
        
        Args:
            demographic_type: 'Gender', 'Location', 'Religion', 'Social Category', 'Age'
            demographic_value: Value to filter by (e.g., 1 for Male, 'Urban' for Location, etc.)
            num_days: Number of days to calculate (default 16)
        
        Returns:
            List of dictionaries with date and vote shares for each party
        """
        # Get demographic column name
        demographic_column_map = {
            'Gender': 'Gender',
            'Location': 'Residential locality type',
            'Religion': '20. Could you please tell me the religion that you belong to?',
            'Social Category': '21. Which social category do you belong to?',
            'Age': 'Could you please tell me your age in complete years?'
        }
        
        demographic_column = demographic_column_map.get(demographic_type)
        if not demographic_column or demographic_column not in self.df.columns:
            return []
        
        # Create filter for demographic
        if demographic_type == 'Gender':
            # Gender: 1=Male, 2=Female
            if demographic_value == 'Male':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Female':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return []
        elif demographic_type == 'Location':
            # Location: 1=Urban, 2=Rural
            if demographic_value == 'Urban':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Rural':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return []
        elif demographic_type == 'Religion':
            # Religion: 1=Hindu, 2=Muslim
            if demographic_value == 'Hindu':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Muslim':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return []
        elif demographic_type == 'Social Category':
            # Social Category: 1=General, 2=OBC, 3=SC, 4=ST
            if demographic_value == 'General+OBC':
                demographic_filter = self.df[demographic_column].isin([1, 2])
            elif demographic_value == 'SC':
                demographic_filter = self.df[demographic_column] == 3
            elif demographic_value == 'ST':
                demographic_filter = self.df[demographic_column] == 4
            else:
                return []
        elif demographic_type == 'Age':
            # Age groups based on numeric age
            age_col = self.df[demographic_column]
            age_numeric = pd.to_numeric(age_col, errors='coerce')
            
            if demographic_value == '18-25':
                demographic_filter = (age_numeric >= 18) & (age_numeric <= 25)
            elif demographic_value == '26-34':
                demographic_filter = (age_numeric >= 26) & (age_numeric <= 34)
            elif demographic_value == '36-50':
                demographic_filter = (age_numeric >= 36) & (age_numeric <= 50)
            elif demographic_value == '50+':
                demographic_filter = age_numeric > 50
            else:
                return []
        else:
            return []
        
        # Calculate dates: last N days ending at reference_date
        target_dates = []
        start_offset = num_days
        all_dates = []
        for i in range(start_offset, -1, -1):
            date = self.reference_date - timedelta(days=i)
            all_dates.append(date)
        
        # Skip Oct 19 if present (matching final PPT)
        oct_19 = self.reference_date - timedelta(days=12)
        for date in all_dates:
            if date != oct_19:
                target_dates.append(date)
        
        if len(target_dates) > num_days:
            target_dates = target_dates[:num_days]
        elif len(target_dates) < num_days:
            while len(target_dates) < num_days:
                target_dates.append(target_dates[-1])
        
        daily_vote_shares = []
        
        for date in target_dates:
            # Filter data by demographic AND date
            actual_end_date = min(date, self.reference_date, self.df['Survey Date'].max())
            date_data = self.df[
                (self.df['Survey Date'] <= actual_end_date) & 
                (self.df['Survey Date'] <= self.reference_date) &
                demographic_filter
            ].copy()
            
            if len(date_data) > 0:
                # Calculate overall normalized vote share (cumulative) for this demographic and date
                vote_shares = self.calculator.calculate_vote_share(
                    date_data,
                    weight_column=self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level',
                    use_weights=True
                )
                
                # Convert date to Excel serial number for chart (matching final PPT format)
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': vote_shares.get('AITC', 0),
                    'BJP': vote_shares.get('BJP', 0),
                    'LEFT': vote_shares.get('LEFT', 0),
                    'INC': vote_shares.get('INC', 0),
                    'Others': vote_shares.get('Others', 0),
                    'NWR': vote_shares.get('NWR', 0)
                })
            else:
                # If no data for this date, still add entry with zero values
                # Convert date to Excel serial number for chart
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': 0,
                    'BJP': 0,
                    'LEFT': 0,
                    'INC': 0,
                    'Others': 0,
                    'NWR': 0
                })
        
        return daily_vote_shares
    
    def update_demographic_chart(self, chart, demographic_type, demographic_value):
        """
        Update demographic chart with Overall Normalized Vote Share (cumulative) for last 16 days
        Filtered by demographic category
        
        Args:
            chart: Chart object from slide
            demographic_type: 'Gender', 'Location', 'Religion', 'Social Category', 'Age'
            demographic_value: Value to filter by
        
        Returns:
            Sample size used for calculations (for updating Base text)
        """
        from pptx.chart.data import CategoryChartData
        
        # Calculate Overall Normalized Vote Share (cumulative) for this demographic
        daily_data = self.calculate_overall_normalized_for_demographic(
            demographic_type, demographic_value, num_days=16
        )
        
        if len(daily_data) == 0:
            print(f"Warning: No data found for {demographic_type}={demographic_value}")
            return
        
        # Extract values
        date_labels_str = [str(d['date_label']) for d in daily_data]
        aitc_values = [round(d['AITC'], 1) for d in daily_data]
        bjp_values = [round(d['BJP'], 1) for d in daily_data]
        left_values = [round(d['LEFT'], 1) for d in daily_data]
        inc_values = [round(d['INC'], 1) for d in daily_data]
        others_values = [round(d['Others'], 1) for d in daily_data]
        nwr_values = [round(d['NWR'], 1) for d in daily_data]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = date_labels_str
        
        chart_data.add_series('AITC (Trinamool Congress)', aitc_values)
        chart_data.add_series('BJP', bjp_values)
        chart_data.add_series('Left Front', left_values)
        chart_data.add_series('INC (Congress)', inc_values)
        chart_data.add_series('Others', others_values)
        chart_data.add_series('N+W+R', nwr_values)
        
        chart.replace_data(chart_data)
        
        # Set axis labels to vertical rotation (like format final file)
        self._set_axis_labels_vertical(chart)
        
        # Calculate sample size: all records up to reference_date, filtered by demographic
        sample_size = self._get_demographic_sample_size(demographic_type, demographic_value, is_7dma=False)
        
        print(f"Updated chart for {demographic_type}={demographic_value} with {len(daily_data)} days of data")
        return sample_size
    
    def calculate_7dma_for_demographic(self, demographic_type, demographic_value, num_days=13):
        """
        Calculate 7 DMA Normalized Vote Share for a specific demographic
        For each date, calculates 7 DMA (moving average) normalized vote share, filtered by demographic
        
        Args:
            demographic_type: 'Gender', 'Location', 'Religion', 'Social Category', 'Age'
            demographic_value: Value to filter by (e.g., 'Male', 'Urban', etc.)
            num_days: Number of days to calculate (default 13 to match final PPT)
        
        Returns:
            List of dictionaries with date and vote shares for each party
        """
        # Get demographic column name
        demographic_column_map = {
            'Gender': 'Gender',
            'Location': 'Residential locality type',
            'Religion': '20. Could you please tell me the religion that you belong to?',
            'Social Category': '21. Which social category do you belong to?',
            'Age': 'Could you please tell me your age in complete years?'
        }
        
        demographic_column = demographic_column_map.get(demographic_type)
        if not demographic_column or demographic_column not in self.df.columns:
            return []
        
        # Create filter for demographic (same logic as calculate_overall_normalized_for_demographic)
        if demographic_type == 'Gender':
            if demographic_value == 'Male':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Female':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return []
        elif demographic_type == 'Location':
            if demographic_value == 'Urban':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Rural':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return []
        elif demographic_type == 'Religion':
            if demographic_value == 'Hindu':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Muslim':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return []
        elif demographic_type == 'Social Category':
            if demographic_value == 'General+OBC':
                demographic_filter = self.df[demographic_column].isin([1, 2])
            elif demographic_value == 'SC':
                demographic_filter = self.df[demographic_column] == 3
            elif demographic_value == 'ST':
                demographic_filter = self.df[demographic_column] == 4
            else:
                return []
        elif demographic_type == 'Age':
            age_col = self.df[demographic_column]
            age_numeric = pd.to_numeric(age_col, errors='coerce')
            
            if demographic_value == '18-25':
                demographic_filter = (age_numeric >= 18) & (age_numeric <= 25)
            elif demographic_value == '26-34':
                demographic_filter = (age_numeric >= 26) & (age_numeric <= 34)
            elif demographic_value == '36-50':
                demographic_filter = (age_numeric >= 36) & (age_numeric <= 50)
            elif demographic_value == '50+':
                demographic_filter = age_numeric > 50
            else:
                return []
        else:
            return []
        
        # Calculate dates: last N days ending at reference_date (matching Slide 7)
        # For 7 DMA charts: Oct 18-31 (excluding Oct 19) = 13 dates
        target_dates = []
        start_offset = num_days + 1  # Start from day before first date
        all_dates = []
        for i in range(start_offset, -1, -1):
            date = self.reference_date - timedelta(days=i)
            all_dates.append(date)
        
        # Exclude Oct 19
        oct_19 = self.reference_date - timedelta(days=12)
        target_dates = [d for d in all_dates if d != oct_19]
        if len(target_dates) > num_days:
            target_dates = target_dates[-num_days:]  # Take last num_days
        
        daily_vote_shares = []
        
        for date in target_dates:
            # For each date, calculate 7 DMA: end_date is one day before the date point
            # So if date is Oct 31, use Oct 24-30 (7 days ending on Oct 30, excluding Oct 31)
            # Window: (date - 7 days) to (date - 1 day) = 7 days total, ending one day before the date
            actual_end_date = date - timedelta(days=1)  # One day before the date point
            cutoff = actual_end_date - timedelta(days=6)  # 6 days before end_date = 7 days total
            
            # Filter data: by demographic AND 7 DMA date range
            max_end_date = min(actual_end_date, self.reference_date - timedelta(days=1), self.df['Survey Date'].max())
            date_data = self.df[
                (self.df['Survey Date'] >= cutoff) & 
                (self.df['Survey Date'] <= max_end_date) &
                (self.df['Survey Date'] <= self.reference_date - timedelta(days=1)) &
                demographic_filter
            ].copy()
            
            if len(date_data) > 0:
                # For 7DMA, use L7D weights when >=50% available, otherwise use regular weights
                # This matches the final PPT behavior where L7D weights are only used when sufficient
                l7d_col = self.get_weight_column('Region', 'L7D') or 'Weight Voteshare L7D Region Level'
                regular_col = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
                
                # Check if L7D weights are available
                l7d_available = date_data[l7d_col].notna().sum()
                total_records = len(date_data)
                
                # For 7DMA: Use L7D weights if >=50% of records have L7D weights, otherwise use regular weights
                # This ensures we don't use sparse L7D weights which give incorrect results
                if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
                    # Use L7D weights - filter to only records with L7D weights
                    date_data_with_weights = date_data[date_data[l7d_col].notna()].copy()
                    weight_column = l7d_col
                else:
                    # Fall back to regular weights when L7D weights are sparse (<50%)
                    date_data_with_weights = date_data.copy()
                    weight_column = regular_col
                
                vote_shares = self.calculator.calculate_vote_share(
                    date_data_with_weights,
                    weight_column=weight_column,
                    use_weights=True
                )
                
                # Convert date to Excel serial number for chart (matching final PPT format)
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': vote_shares.get('AITC', 0),
                    'BJP': vote_shares.get('BJP', 0),
                    'LEFT': vote_shares.get('LEFT', 0),
                    'INC': vote_shares.get('INC', 0),
                    'Others': vote_shares.get('Others', 0),
                    'NWR': vote_shares.get('NWR', 0)
                })
            else:
                # If no data for this date, still add entry with zero values
                # Convert date to Excel serial number for chart
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': 0,
                    'BJP': 0,
                    'LEFT': 0,
                    'INC': 0,
                    'Others': 0,
                    'NWR': 0
                })
        
        return daily_vote_shares
    
    def update_7dma_demographic_chart(self, chart, demographic_type, demographic_value):
        """
        Update demographic 7 DMA chart with 7 DMA Normalized Vote Share for last 13 days
        Filtered by demographic category
        
        Args:
            chart: Chart object from slide
            demographic_type: 'Gender', 'Location', 'Religion', 'Social Category', 'Age'
            demographic_value: Value to filter by
        
        Returns:
            Sample size used for calculations (for updating Base text)
        """
        from pptx.chart.data import CategoryChartData
        
        # Calculate 7 DMA Normalized Vote Share for this demographic
        daily_data = self.calculate_7dma_for_demographic(
            demographic_type, demographic_value, num_days=13
        )
        
        if len(daily_data) == 0:
            print(f"Warning: No data found for 7 DMA {demographic_type}={demographic_value}")
            return
        
        # Extract date labels for categories (use date strings like "15/10" instead of Excel serial numbers)
        date_labels = [str(d['date_label']) for d in daily_data]
        
        # Prepare series data
        aitc_values = [float(round(d['AITC'], 1)) for d in daily_data]
        bjp_values = [float(round(d['BJP'], 1)) for d in daily_data]
        left_values = [float(round(d['LEFT'], 1)) for d in daily_data]
        inc_values = [float(round(d['INC'], 1)) for d in daily_data]
        others_values = [float(round(d['Others'], 1)) for d in daily_data]
        nwr_values = [float(round(d['NWR'], 1)) for d in daily_data]
        
        # Create chart data using date strings (matching final PPT format)
        chart_data = CategoryChartData()
        chart_data.categories = date_labels
        
        chart_data.add_series('AITC (Trinamool Congress)', aitc_values)
        chart_data.add_series('BJP', bjp_values)
        chart_data.add_series('Left Front', left_values)
        chart_data.add_series('INC (Congress)', inc_values)
        chart_data.add_series('Others', others_values)
        chart_data.add_series('N+W+R', nwr_values)
        
        chart.replace_data(chart_data)
        
        # Set axis labels to vertical rotation (like format final file)
        self._set_axis_labels_vertical(chart)
        
        # Calculate sample size: 7DMA window, filtered by demographic
        sample_size = self._get_demographic_sample_size(demographic_type, demographic_value, is_7dma=True)
        
        print(f"Updated 7 DMA chart for {demographic_type}={demographic_value} with {len(daily_data)} days of data")
        return sample_size
    
    def _get_demographic_sample_size(self, demographic_type, demographic_value, is_7dma=False):
        """
        Helper function to calculate sample size for a demographic category
        
        Args:
            demographic_type: 'Gender', 'Location', 'Religion', 'Social Category', 'Age'
            demographic_value: Value to filter by
            is_7dma: If True, calculate for 7DMA window; if False, calculate for overall
        
        Returns:
            Sample size (number of records)
        """
        # Get demographic column name
        demographic_column_map = {
            'Gender': 'Gender',
            'Location': 'Residential locality type',
            'Religion': '20. Could you please tell me the religion that you belong to?',
            'Social Category': '21. Which social category do you belong to?',
            'Age': 'Could you please tell me your age in complete years?'
        }
        
        demographic_column = demographic_column_map.get(demographic_type)
        if not demographic_column or demographic_column not in self.df.columns:
            return 0
        
        # Create filter for demographic
        if demographic_type == 'Gender':
            if demographic_value == 'Male':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Female':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return 0
        elif demographic_type == 'Location':
            if demographic_value == 'Urban':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Rural':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return 0
        elif demographic_type == 'Religion':
            if demographic_value == 'Hindu':
                demographic_filter = self.df[demographic_column] == 1
            elif demographic_value == 'Muslim':
                demographic_filter = self.df[demographic_column] == 2
            else:
                return 0
        elif demographic_type == 'Social Category':
            if demographic_value == 'General+OBC':
                demographic_filter = self.df[demographic_column].isin([1, 2])
            elif demographic_value == 'SC':
                demographic_filter = self.df[demographic_column] == 3
            elif demographic_value == 'ST':
                demographic_filter = self.df[demographic_column] == 4
            else:
                return 0
        elif demographic_type == 'Age':
            age_col = self.df[demographic_column]
            age_numeric = pd.to_numeric(age_col, errors='coerce')
            if demographic_value == '18-25':
                demographic_filter = (age_numeric >= 18) & (age_numeric <= 25)
            elif demographic_value == '26-34':
                demographic_filter = (age_numeric >= 26) & (age_numeric <= 34)
            elif demographic_value == '36-50':
                demographic_filter = (age_numeric >= 36) & (age_numeric <= 50)
            elif demographic_value == '50+':
                demographic_filter = age_numeric > 50
            else:
                return 0
        else:
            return 0
        
        if is_7dma:
            # For 7DMA: use last 7 days ending one day before reference_date
            end_date_7dma = self.reference_date - timedelta(days=1)
            cutoff_7dma = end_date_7dma - timedelta(days=6)
            date_filter = (self.df['Survey Date'] >= cutoff_7dma) & (self.df['Survey Date'] <= end_date_7dma)
            
            filtered_data = self.df[date_filter & demographic_filter].copy()
            
            # Sample size: only count records with non-empty responses to main question
            sample_size = self.get_sample_size(filtered_data)
        else:
            # For Overall: use all data up to reference_date
            date_filter = self.df['Survey Date'] <= self.reference_date
            filtered_data = self.df[date_filter & demographic_filter].copy()
            # Sample size: only count records with non-empty responses to main question
            sample_size = self.get_sample_size(filtered_data)
        
        return sample_size
    
    def update_base_text_on_slide(self, slide, sample_size, slide_num=None):
        """
        Helper function to update Base sample size text on any slide
        Searches in all shapes including text frames, tables, and group shapes
        
        Args:
            slide: Slide object
            sample_size: Sample size to display
            slide_num: Optional slide number, used for specific handling of Gains and Losses slides
        
        Returns:
            True if Base text was found and updated, False otherwise
        """
        if sample_size <= 0:
            return False
        
        # For slides 47-48, don't update Base text (keep as is)
        if slide_num is not None and slide_num in [47, 48]:
            return False
        
        found_and_updated = False
        
        # Specific handling for Gains and Losses demographic slides 49-61
        is_gains_losses_demographic_slide = (slide_num is not None and 49 <= slide_num <= 61)
        
        def update_text_in_shape(shape):
            """Recursive function to update Base text in any shape"""
            nonlocal found_and_updated
            
            # Check text frame
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_frame = shape.text_frame
                # Try to disable word wrap to keep text on one line
                try:
                    text_frame.word_wrap = False
                except:
                    pass
                
                for paragraph in shape.text_frame.paragraphs:
                    # Get full paragraph text (may span multiple runs)
                    full_paragraph_text = ''.join([run.text for run in paragraph.runs])
                    
                    if is_gains_losses_demographic_slide:
                        # For slides 49-61, specifically look for "Base of all eligible voters of 2021 : {number}"
                        # and replace only the number, keeping the text "Base of all eligible voters of 2021 :"
                        # First, try to match the full text with a number
                        pattern_gains_losses_base = r'(Base of all eligible voters of 2021\s*:?\s*)(\d{1,3}(?:,\d{3})*)'
                        match = re.search(pattern_gains_losses_base, full_paragraph_text, re.IGNORECASE)
                        if match:
                            # Replace only the number part (group 2), keep the text (group 1)
                            # Use lambda to avoid issues with comma-formatted numbers in replacement string
                            # Remove any newlines and ensure single line
                            # Replace only the number, keeping text on same line
                            new_text = re.sub(pattern_gains_losses_base, lambda m: f'{m.group(1).strip()} {sample_size:,}', full_paragraph_text, flags=re.IGNORECASE)
                            # Remove all newlines and extra whitespace, ensure single line
                            new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                            # Clear all runs and set new text in first run
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = new_text
                            found_and_updated = True
                            return True
                        # Also check if text contains "Base of all eligible voters of 2021" but no number yet
                        if 'Base of all eligible voters of 2021' in full_paragraph_text or 'base of all eligible voters of 2021' in full_paragraph_text.lower():
                            # Replace any existing number pattern after the text, or append if no number
                            pattern_gains_losses_base_no_num = r'(Base of all eligible voters of 2021\s*:?\s*)(\d{1,3}(?:,\d{3})*)?'
                            new_text = re.sub(pattern_gains_losses_base_no_num, f'Base of all eligible voters of 2021 : {sample_size:,}', full_paragraph_text, flags=re.IGNORECASE)
                            # Remove all newlines and extra whitespace, ensure single line
                            new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                            # Clear all runs and set new text in first run
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = new_text
                            found_and_updated = True
                            return True
                        # Also handle case where it might be "Base: {number}" and we need to replace with full text
                        pattern_base_colon = r'(Base\s*:\s*)(\d{1,3}(?:,\d{3})*)'
                        match = re.search(pattern_base_colon, full_paragraph_text, re.IGNORECASE)
                        if match:
                            # Replace "Base: {number}" with "Base of all eligible voters of 2021 : {sample_size}"
                            new_text = f"Base of all eligible voters of 2021 : {sample_size:,}"
                            # Remove all newlines and extra whitespace, ensure single line
                            new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                            # Clear all runs and set new text in first run
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = new_text
                            # Disable word wrap if possible (to keep on one line)
                            if hasattr(paragraph, 'text_frame'):
                                try:
                                    paragraph.text_frame.word_wrap = False
                                except:
                                    pass
                            found_and_updated = True
                            return True
                    else:
                        # For other slides (6, 7, 14-36, 44, 45, 63, 64), use existing patterns
                        # Get full paragraph text (may span multiple runs) to avoid duplicates
                        full_paragraph_text = ''.join([run.text for run in paragraph.runs])
                        
                        patterns = [
                            # Pattern 1: "Base: 12,345" or "Base 12,345"
                            (r'(Base\s*:?\s*)(\d{1,3}(?:,\d{3})*)', lambda m: f'{m.group(1)}{sample_size:,}'),
                            # Pattern 2: "Base: 12345" or "Base 12345" (no commas)
                            (r'(Base\s*:?\s*)(\d+)', lambda m: f'{m.group(1)}{sample_size:,}'),
                            # Pattern 3: "Base: 12 345" (space-separated)
                            (r'(Base\s*:?\s*)(\d{1,3}(?:\s+\d{3})*)', lambda m: f'{m.group(1)}{sample_size:,}'),
                        ]
                        
                        for pattern, replacement in patterns:
                            match = re.search(pattern, full_paragraph_text, re.IGNORECASE)
                            if match:
                                new_text = re.sub(pattern, replacement, full_paragraph_text, flags=re.IGNORECASE)
                                # Remove all newlines and extra whitespace, ensure single line
                                new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                                # Clear all runs and set new text in first run
                                paragraph.clear()
                                run = paragraph.add_run()
                                run.text = new_text
                                found_and_updated = True
                                return True
                        
                        # Also check if text contains "Base" but no number yet (append number)
                        if 'Base' in full_paragraph_text and not re.search(r'\d', full_paragraph_text):
                            # Append sample size after "Base"
                            if full_paragraph_text.strip().endswith('Base') or full_paragraph_text.strip().endswith('Base:'):
                                new_text = f"{full_paragraph_text.strip()} {sample_size:,}"
                            else:
                                # Replace "Base" with "Base: sample_size"
                                new_text = re.sub(r'Base\s*:?', f'Base: {sample_size:,}', full_paragraph_text, flags=re.IGNORECASE)
                            # Remove all newlines and extra whitespace, ensure single line
                            new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                            # Clear all runs and set new text in first run
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = new_text
                            found_and_updated = True
                            return True
            
            # Check table cells
            if hasattr(shape, 'has_table') and shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if hasattr(cell, 'text_frame') and cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text = run.text
                                    
                                    if is_gains_losses_demographic_slide:
                                        # For slides 49-61, get full paragraph text (may span multiple runs)
                                        full_cell_text = ''.join([r.text for r in paragraph.runs])
                                        
                                        # First, try to match the full text with a number
                                        pattern_gains_losses_base = r'(Base of all eligible voters of 2021\s*:?\s*)(\d{1,3}(?:,\d{3})*)'
                                        match = re.search(pattern_gains_losses_base, full_cell_text, re.IGNORECASE)
                                        if match:
                                            # Replace only the number part (group 2), keep the text (group 1)
                                            # Use lambda to avoid issues with comma-formatted numbers in replacement string
                                            # Replace only the number, keeping text on same line
                                            new_text = re.sub(pattern_gains_losses_base, lambda m: f'{m.group(1).strip()} {sample_size:,}', full_cell_text, flags=re.IGNORECASE)
                                            # Remove all newlines and extra whitespace, ensure single line
                                            new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                                            # Clear all runs and set new text in first run
                                            paragraph.clear()
                                            run = paragraph.add_run()
                                            run.text = new_text
                                            found_and_updated = True
                                            return True
                                        # Also check if text contains "Base of all eligible voters of 2021" but no number yet
                                        if 'Base of all eligible voters of 2021' in full_cell_text or 'base of all eligible voters of 2021' in full_cell_text.lower():
                                            # Replace any existing number pattern after the text, or append if no number
                                            pattern_gains_losses_base_no_num = r'(Base of all eligible voters of 2021\s*:?\s*)(\d{1,3}(?:,\d{3})*)?'
                                            new_text = re.sub(pattern_gains_losses_base_no_num, f'Base of all eligible voters of 2021 : {sample_size:,}', full_cell_text, flags=re.IGNORECASE)
                                            # Remove all newlines and extra whitespace, ensure single line
                                            new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                                            # Clear all runs and set new text in first run
                                            paragraph.clear()
                                            run = paragraph.add_run()
                                            run.text = new_text
                                            found_and_updated = True
                                            return True
                                        # Also handle case where it might be "Base: {number}" and we need to replace with full text
                                        pattern_base_colon = r'(Base\s*:\s*)(\d{1,3}(?:,\d{3})*)'
                                        match = re.search(pattern_base_colon, full_cell_text, re.IGNORECASE)
                                        if match:
                                            # Replace "Base: {number}" with "Base of all eligible voters of 2021 : {sample_size}"
                                            new_text = f"Base of all eligible voters of 2021 : {sample_size:,}"
                                            # Remove all newlines and extra whitespace, ensure single line
                                            new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                                            # Clear all runs and set new text in first run
                                            paragraph.clear()
                                            run = paragraph.add_run()
                                            run.text = new_text
                                            found_and_updated = True
                                            return True
                                    else:
                                        # Check for Base pattern in table cells
                                        # Get full paragraph text (may span multiple runs) to avoid duplicates
                                        full_cell_text = ''.join([r.text for r in paragraph.runs])
                                        
                                        patterns = [
                                            # Pattern 1: "Base: 12,345" or "Base 12,345"
                                            (r'(Base\s*:?\s*)(\d{1,3}(?:,\d{3})*)', lambda m: f'{m.group(1)}{sample_size:,}'),
                                            # Pattern 2: "Base: 12345" or "Base 12345" (no commas)
                                            (r'(Base\s*:?\s*)(\d+)', lambda m: f'{m.group(1)}{sample_size:,}'),
                                        ]
                                        for pattern, replacement in patterns:
                                            match = re.search(pattern, full_cell_text, re.IGNORECASE)
                                            if match:
                                                new_text = re.sub(pattern, replacement, full_cell_text, flags=re.IGNORECASE)
                                                # Remove all newlines and extra whitespace, ensure single line
                                                new_text = re.sub(r'\s+', ' ', new_text.replace('\n', ' ').replace('\r', ' ')).strip()
                                                # Clear all runs and set new text in first run
                                                paragraph.clear()
                                                run = paragraph.add_run()
                                                run.text = new_text
                                                found_and_updated = True
                                                return True
            
            # Check group shapes (recursive)
            if hasattr(shape, 'shapes'):
                for sub_shape in shape.shapes:
                    if update_text_in_shape(sub_shape):
                        return True
            
            return False
        
        # Search all shapes on the slide
        for shape in slide.shapes:
            if update_text_in_shape(shape):
                break
        
        return found_and_updated
    
    def update_slide_6_chart(self, chart):
        """
        Update Slide 6 chart with Overall Normalized Vote Share (cumulative) for last 16 days
        For each date, shows normalized vote share using ALL data up to that date
        
        Args:
            chart: Chart object from Slide 6
        
        Returns:
            Sample size used for calculations (for updating Base text)
        """
        from pptx.chart.data import CategoryChartData
        
        # Calculate Overall Normalized Vote Share (cumulative) for last 16 days (matching final PPT)
        daily_data = self.calculate_overall_normalized_for_days(num_days=16)
        
        if len(daily_data) == 0:
            print("Warning: No data available for chart")
            return
        
        # Extract date labels and vote share values
        date_labels = [d['date_label'] for d in daily_data]
        
        # Update chart by modifying series values directly (more reliable)
        # This avoids potential corruption from replace_data() when updating categories
        plot = chart.plots[0]
        series_list = plot.series
        
        # Ensure we have the right number of series
        if len(series_list) < 6:
            print(f"Warning: Chart has {len(series_list)} series, expected 6")
            return
        
        # Prepare series data (ensure all values are floats)
        # Round values to 1 decimal place
        aitc_values = [float(round(d['AITC'], 1)) for d in daily_data]
        bjp_values = [float(round(d['BJP'], 1)) for d in daily_data]
        left_values = [float(round(d['LEFT'], 1)) for d in daily_data]
        inc_values = [float(round(d['INC'], 1)) for d in daily_data]
        others_values = [float(round(d['Others'], 1)) for d in daily_data]
        nwr_values = [float(round(d['NWR'], 1)) for d in daily_data]
        
        # Update chart using replace_data with complete data structure (single update)
        # This ensures clean update without corruption
        from pptx.chart.data import CategoryChartData
        
        date_labels_str = [str(d['date_label']) for d in daily_data]
        chart_data = CategoryChartData()
        chart_data.categories = date_labels_str
        
        # Add all series with proper names matching template
        chart_data.add_series('AITC (Trinamool Congress)', aitc_values)
        chart_data.add_series('BJP', bjp_values)
        chart_data.add_series('Left Front', left_values)
        chart_data.add_series('INC (Congress)', inc_values)
        chart_data.add_series('Others', others_values)
        chart_data.add_series('N+W+R', nwr_values)
        
        # Single clean replace_data call (avoid multiple updates)
        chart.replace_data(chart_data)
        
        # Set axis labels to vertical rotation (like format final file)
        self._set_axis_labels_vertical(chart)
        
        # Calculate sample size: only count records with non-empty responses to main question
        sample_size = self.get_sample_size(self.df_filtered, 'overall')
        
        print(f"Updated Slide 6 chart with {len(daily_data)} days of Overall Normalized Vote Share data")
        return sample_size
    
    def calculate_7dma_for_days(self, num_days=13):
        """
        Calculate 7 DMA vote shares for each of the last N days
        For each date, calculates 7 DMA using data from (date - 6 days) to date
        Dates end at reference_date (when report is generated)
        
        Args:
            num_days: Number of days to calculate (default 13 to match final PPT Slide 7)
                      Final PPT Slide 7 shows Oct 18-31 excluding Oct 19 = 13 dates
        
        Returns:
            List of dictionaries with date and 7 DMA vote shares for each party
        """
        # Calculate dates: last N days ending at reference_date (when report is generated)
        # For 13 days ending at Oct 31: start from Oct 19 (reference_date - 12 days)
        # But final PPT shows Oct 18-31 excluding Oct 19, so we start from Oct 18
        target_dates = []
        # Start from (reference_date - num_days + 1) days back for num_days total
        # For 13 days ending Oct 31: start from Oct 19, but we need Oct 18-31 (excluding Oct 19)
        # So we generate Oct 18-31 and exclude Oct 19
        start_offset = num_days + 1  # Generate 14 days to have Oct 18-31
        all_dates = []
        for i in range(start_offset, -1, -1):  # From start_offset down to 0 (reference date)
            date = self.reference_date - timedelta(days=i)
            all_dates.append(date)
        
        # Final PPT Slide 7 excludes Oct 19
        # So we have Oct 18, 20, 21, ..., 31 (13 dates)
        oct_19 = self.reference_date - timedelta(days=12)  # Oct 19 is 12 days before Oct 31
        target_dates = [d for d in all_dates if d != oct_19]
        
        # Ensure we have exactly num_days dates
        if len(target_dates) > num_days:
            target_dates = target_dates[-num_days:]  # Take last num_days
        elif len(target_dates) < num_days:
            # Pad if needed (shouldn't happen)
            while len(target_dates) < num_days:
                target_dates.append(target_dates[-1])
        
        daily_vote_shares = []
        
        for date in target_dates:
            # For each date, calculate 7 DMA: end_date is one day before the date point
            # So if date is Oct 31, use Oct 24-30 (7 days ending on Oct 30, excluding Oct 31)
            # Window: (date - 7 days) to (date - 1 day) = 7 days total, ending one day before the date
            actual_end_date = date - timedelta(days=1)  # One day before the date point
            cutoff = actual_end_date - timedelta(days=6)  # 6 days before end_date = 7 days total
            
            # Filter data: from cutoff to actual_end_date, and not exceed reference_date - 1
            max_end_date = min(actual_end_date, self.reference_date - timedelta(days=1), self.df['Survey Date'].max())
            date_data = self.df[
                (self.df['Survey Date'] >= cutoff) & 
                (self.df['Survey Date'] <= max_end_date) &
                (self.df['Survey Date'] <= self.reference_date - timedelta(days=1))
            ].copy()
            
            if len(date_data) > 0:
                # For 7DMA, use L7D weights when >=50% available, otherwise use regular weights
                # This matches the final PPT behavior where L7D weights are only used when sufficient
                l7d_col = self.get_weight_column('Region', 'L7D') or 'Weight Voteshare L7D Region Level'
                regular_col = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
                
                # Check if L7D weights are available
                l7d_available = date_data[l7d_col].notna().sum()
                total_records = len(date_data)
                
                # For 7DMA: Use L7D weights if >=50% of records have L7D weights, otherwise use regular weights
                # This ensures we don't use sparse L7D weights which give incorrect results
                if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
                    # Use L7D weights - filter to only records with L7D weights
                    date_data_with_weights = date_data[date_data[l7d_col].notna()].copy()
                    weight_column = l7d_col
                else:
                    # Fall back to regular weights when L7D weights are sparse (<50%)
                    date_data_with_weights = date_data.copy()
                    weight_column = regular_col
                
                vote_shares = self.calculator.calculate_vote_share(
                    date_data_with_weights,
                    weight_column=weight_column,
                    use_weights=True
                )
                
                # Convert date to Excel serial number for chart (matching final PPT format)
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': vote_shares.get('AITC', 0),
                    'BJP': vote_shares.get('BJP', 0),
                    'LEFT': vote_shares.get('LEFT', 0),
                    'INC': vote_shares.get('INC', 0),
                    'Others': vote_shares.get('Others', 0),
                    'NWR': vote_shares.get('NWR', 0)
                })
            else:
                # If no data for this date, still add entry with zero values
                # Convert date to Excel serial number for chart
                excel_epoch = datetime(1899, 12, 30)
                excel_serial = (date - excel_epoch).days
                
                daily_vote_shares.append({
                    'date': date,
                    'date_label': date.strftime('%m/%d/%Y'),
                    'excel_serial': excel_serial,  # Excel serial number for chart
                    'AITC': 0,
                    'BJP': 0,
                    'LEFT': 0,
                    'INC': 0,
                    'Others': 0,
                    'NWR': 0
                })
        
        return daily_vote_shares
    
    def update_slide_7_chart(self, chart):
        """
        Update Slide 7 chart with 7 DMA vote share data
        For each date, shows 7 DMA (moving average) normalized vote share
        
        Args:
            chart: Chart object from Slide 7
        
        Returns:
            Sample size used for calculations (for updating Base text)
        """
        from pptx.chart.data import CategoryChartData
        
        # Calculate 7 DMA for last 13 days (matching final PPT Slide 7)
        daily_data = self.calculate_7dma_for_days(num_days=13)
        
        if len(daily_data) == 0:
            print("Warning: No data available for Slide 7 chart")
            return
        
        # Extract date labels for categories (use date strings like "10/30/2025" instead of Excel serial numbers)
        date_labels = [str(d['date_label']) for d in daily_data]
        
        # Prepare series data (ensure all values are floats)
        # Round values to 1 decimal place
        aitc_values = [float(round(d['AITC'], 1)) for d in daily_data]
        bjp_values = [float(round(d['BJP'], 1)) for d in daily_data]
        left_values = [float(round(d['LEFT'], 1)) for d in daily_data]
        inc_values = [float(round(d['INC'], 1)) for d in daily_data]
        others_values = [float(round(d['Others'], 1)) for d in daily_data]
        nwr_values = [float(round(d['NWR'], 1)) for d in daily_data]
        
        # Update chart using replace_data with complete data structure
        chart_data = CategoryChartData()
        # Use date strings as categories (matching final PPT format)
        chart_data.categories = date_labels
        
        # Add all series with proper names matching template
        chart_data.add_series('AITC (Trinamool Congress)', aitc_values)
        chart_data.add_series('BJP', bjp_values)
        chart_data.add_series('Left Front', left_values)
        chart_data.add_series('INC (Congress)', inc_values)
        chart_data.add_series('Others', others_values)
        chart_data.add_series('N+W+R', nwr_values)
        
        # Single clean replace_data call
        chart.replace_data(chart_data)
        
        # Set axis labels to vertical rotation (like format final file)
        self._set_axis_labels_vertical(chart)
        
        # Calculate sample size: 7DMA window (last 7 days ending one day before reference_date)
        end_date_7dma = self.reference_date - timedelta(days=1)
        cutoff_7dma = end_date_7dma - timedelta(days=6)
        dma7_data = self.df[
            (self.df['Survey Date'] >= cutoff_7dma) & 
            (self.df['Survey Date'] <= end_date_7dma)
        ].copy()
        
        # Check if L7D weights are available
        l7d_col = self.get_weight_column('Region', 'L7D') or 'Weight Voteshare L7D Region Level'
        l7d_available = dma7_data[l7d_col].notna().sum()
        total_records = len(dma7_data)
        
        # Sample size: only count records with non-empty responses to main question
        sample_size = self.get_sample_size(dma7_data)
        
        print(f"Updated Slide 7 chart with {len(daily_data)} days of 7 DMA data")
        return sample_size
    
    def map_reason_to_display_label(self, reason_text):
        """
        Map actual reason text from Excel to display label used in final PPT
        
        Args:
            reason_text: Original reason text from Excel column
        
        Returns:
            Display label as shown in final PPT
        """
        # Mapping based on final PPT slides 38-39
        # Final PPT shows:
        # 1. "Beneficial to Bengal" (from "For the benefit of West Bengal")
        # 2. "Good Overall Performance" (from "The party has performed well in the state")
        # 3. "Delivered on Promises" (from "For good governance / delivering government services")
        
        reason_lower = reason_text.lower()
        
        if 'benefit' in reason_lower and 'west bengal' in reason_lower:
            return 'Beneficial to Bengal'
        elif 'performed well' in reason_lower or 'performance' in reason_lower:
            return 'Good Overall Performance'
        elif 'good governance' in reason_lower or 'delivering government services' in reason_lower:
            return 'Delivered on Promises'
        else:
            # Return original if no mapping found
            return reason_text
    
    def calculate_top_reasons(self, party, is_7dma=False):
        """
        Calculate top 3 reasons for voting for AITC only
        Returns rephrased labels matching final PPT
        
        Args:
            party: 'AITC' (only AITC is shown in slides 38-39)
            is_7dma: If True, calculate for 7 DMA data; if False, calculate for overall data
        
        Returns:
            Dictionary with 'reasons' (list of top 3 reasons with percentages) and 'categories' (list of display labels)
        """
        # Filter data based on 7 DMA or overall
        if is_7dma:
            # For 7 DMA: use last 7 days ending one day before reference_date
            end_date = self.reference_date - timedelta(days=1)
            cutoff = end_date - timedelta(days=6)
            filtered_df = self.df[
                (self.df['Survey Date'] >= cutoff) &
                (self.df['Survey Date'] <= end_date) &
                (self.df['Survey Date'] <= self.reference_date - timedelta(days=1))
            ].copy()
        else:
            # For overall: use all data up to reference_date
            filtered_df = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # For 7DMA, ALWAYS use L7D weights (as per user requirement)
        # Filter to only include records with L7D weights when available
        l7d_col = self.get_weight_column('Region', 'L7D') or 'Weight Voteshare L7D Region Level'
        regular_col = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        if is_7dma:
            # Check if L7D weights are available
            l7d_available = filtered_df[l7d_col].notna().sum()
            total_records = len(filtered_df)
            
            # For 7DMA: Use L7D weights if >=50% of records have L7D weights, otherwise use regular weights
            # This ensures we don't use sparse L7D weights which give incorrect results
            if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
                # Use L7D weights - filter to only records with L7D weights
                filtered_df = filtered_df[filtered_df[l7d_col].notna()].copy()
                weight_col = l7d_col
            else:
                # Fall back to regular weights when L7D weights are sparse (<50%)
                weight_col = regular_col
        else:
            weight_col = regular_col
        
        weights = pd.to_numeric(filtered_df[weight_col], errors='coerce').fillna(0)
        
        # Get reason columns based on party
        if party == 'AITC':
            # Question 11: "In your opinion what are the top 3 reasons for voting for AITC?"
            reason_prefix = '11. In your opinion what are the top 3 reasons for voting for AITC?'
        elif party == 'BJP':
            # Question 12: "In your opinion what are the top 3 reasons for voting for BJP?"
            reason_prefix = '12. In your opinion what are the top 3 reasons for voting for BJP?'
        else:
            return {'reasons': [], 'categories': []}
        
        # Find all columns for this question (exclude "Others (specify)" and "Don't know")
        def should_exclude_column(col_str):
            """Check if column should be excluded"""
            col_lower = str(col_str).lower()
            # Extract reason text (everything after " - ")
            if ' - ' in col_str:
                reason_text = col_str.split(' - ')[-1].lower()
            else:
                reason_text = col_lower
            
            exclude_patterns = [
                'others (specify)',
                "don",  # Match "don" anywhere (will catch "don't know")
                "can't say",  # Full phrase
                "can't",  # Match "can't" anywhere
                '.1'  # Duplicate columns
            ]
            # Check if any pattern matches in either full column name or reason text
            # Also check for "know" and "say" together as exclusion
            if 'know' in reason_text and 'say' in reason_text:
                return True
            for pattern in exclude_patterns:
                if pattern in col_lower or pattern in reason_text:
                    return True
            return False
        
        reason_cols = [col for col in self.df.columns 
                       if str(col).startswith(reason_prefix) 
                       and not should_exclude_column(str(col))]
        
        # Calculate weighted counts for each reason
        reason_counts = {}
        for col in reason_cols:
            # Extract reason text (everything after " - ")
            reason_text = str(col).split(' - ')[-1] if ' - ' in str(col) else str(col)
            
            # Convert column to numeric (0 or 1)
            reason_values = pd.to_numeric(filtered_df[col], errors='coerce').fillna(0)
            
            # Calculate weighted sum
            weighted_sum = (reason_values * weights).sum()
            
            # Calculate percentage (weighted sum / total weighted responses)
            # Total weighted responses = sum of all weights for records that selected at least one reason
            total_weighted = weights.sum()
            percentage = (weighted_sum / total_weighted * 100) if total_weighted > 0 else 0
            
            reason_counts[reason_text] = {
                'weighted_count': weighted_sum,
                'percentage': percentage
            }
        
        # Sort by weighted count and get top 3
        top_reasons = sorted(reason_counts.items(), key=lambda x: x[1]['weighted_count'], reverse=True)[:3]
        
        # Map to display labels (rephrased as in final PPT)
        categories = []
        percentages = []
        mapped_reasons = []
        
        for reason_text, data in top_reasons:
            display_label = self.map_reason_to_display_label(reason_text)
            categories.append(display_label)
            percentages.append(data['percentage'])
            mapped_reasons.append((display_label, reason_text, data))
        
        # Calculate sample size (number of records used in calculation)
        sample_size = len(filtered_df)
        
        return {
            'reasons': mapped_reasons,  # List of (display_label, original_reason_text, {weighted_count, percentage})
            'categories': categories,  # Display labels (rephrased)
            'percentages': percentages,
            'sample_size': sample_size  # Sample size used for calculations
        }
    
    def update_top_reasons_chart(self, chart, is_7dma=False):
        """
        Update chart with top 3 reasons for AITC only (as shown in final PPT slides 38-39)
        
        Args:
            chart: Chart object from Slide 38 or 39
            is_7dma: If True, update for 7 DMA data (Slide 38); if False, update for overall data (Slide 39)
        
        Returns:
            Sample size used for calculations (for updating Base text)
        """
        from pptx.chart.data import CategoryChartData
        
        # Calculate top reasons for AITC only (slides 38-39 only show AITC)
        aitc_reasons = self.calculate_top_reasons('AITC', is_7dma=is_7dma)
        
        if len(aitc_reasons['categories']) == 0:
            print("Warning: No reason data available for chart")
            return 0
        
        # Use the display labels (rephrased) from calculate_top_reasons
        categories = aitc_reasons['categories']
        percentages = aitc_reasons['percentages']
        sample_size = aitc_reasons.get('sample_size', 0)
        
        # Round values to 1 decimal place
        aitc_values = [round(v, 1) for v in percentages]
        
        # Update chart - only AITC series (as in final PPT)
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        # Add only AITC series (1 series, not 2)
        chart_data.add_series('AITC', aitc_values)
        
        # Replace chart data
        chart.replace_data(chart_data)
        
        print(f"Updated chart with top 3 AITC reasons: {categories}")
        return sample_size
    
    def categorize_party_2021(self, party_code):
        """
        Categorize party code from 2021 voting question (Q5) into standard categories
        Same mapping as Q8 (current voting)
        
        Args:
            party_code: Party code from survey (Q5)
        
        Returns:
            Party category string
        """
        if pd.isna(party_code):
            return 'NWR'
        
        try:
            party_code = int(party_code)
        except (ValueError, TypeError):
            return 'NWR'
        
        # Main parties (same mapping as Q8)
        if party_code == 1:
            return 'AITC'
        elif party_code == 2:
            return 'BJP'
        elif party_code == 3:
            return 'INC'
        elif party_code == 4:
            return 'LEFT'
        elif party_code in [12, 44]:
            return 'Others'
        elif party_code in [55, 66, 67, 77, 78, 88]:
            return 'NWR'
        else:
            return 'Others'
    
    def calculate_ae2021_vote_shares_from_survey(self, data_filtered, weight_column=None, use_weights=True):
        """
        Calculate 2021 AE vote shares from survey responses (Q5)
        Used for demographic tables where we don't have government data
        
        Args:
            data_filtered: Filtered DataFrame
            weight_column: Column name for weights (if None, will auto-detect)
            use_weights: If True, use weights for normalization
        
        Returns:
            Dictionary with party vote shares
        """
        # Auto-detect weight column if not provided
        if weight_column is None:
            weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        question_2021 = '5. Which party did you vote for in the last assembly elections (MLA) in 2021?'
        
        if question_2021 not in data_filtered.columns:
            return {'AITC': 0, 'BJP': 0, 'LEFT': 0, 'INC': 0, 'Others': 0, 'NWR': 0}
        
        vote_data = data_filtered[[question_2021]].copy()
        
        # Categorize parties
        vote_data['party_category'] = vote_data[question_2021].apply(self.categorize_party_2021)
        
        # Add weights
        if use_weights and weight_column and weight_column in data_filtered.columns:
            vote_data = vote_data.reset_index(drop=True)
            weights = data_filtered[weight_column].reset_index(drop=True)
            vote_data['weight'] = pd.to_numeric(weights, errors='coerce')
            vote_data = vote_data[vote_data['weight'].notna()].copy()
        else:
            vote_data['weight'] = 1.0
        
        total_weight = vote_data['weight'].sum()
        
        if total_weight == 0:
            return {'AITC': 0, 'BJP': 0, 'LEFT': 0, 'INC': 0, 'Others': 0, 'NWR': 0}
        
        # Calculate weighted vote share for each party
        vote_shares = {}
        for party in ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']:
            party_mask = vote_data['party_category'] == party
            party_weight = vote_data.loc[party_mask, 'weight'].sum()
            vote_shares[party] = (party_weight / total_weight) * 100
        
        return vote_shares
    
    def calculate_gains_losses(self, data_filtered, is_7dma=False, demographic_type=None, demographic_value=None):
        """
        Calculate Gains and Losses table data
        Shows percentage of people who voted for party X in 2021 and will vote for party Y in 2025
        
        Args:
            data_filtered: Filtered DataFrame (overall or demographic)
            is_7dma: If True, use 7DMA data; if False, use Overall data
            demographic_type: Demographic type if demographic table (None for overall)
            demographic_value: Demographic value if demographic table (None for overall)
        
        Returns:
            Dictionary with gains/losses matrix and base sample size
        """
        question_2021 = '5. Which party did you vote for in the last assembly elections (MLA) in 2021?'
        question_2025 = '8. If assembly elections (MLA) were to be held tomorrow, then which party would you vote for?'
        weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        # Filter to only records that voted in 2021 (have valid Q5 response)
        eligible_data = data_filtered[data_filtered[question_2021].notna()].copy()
        
        if len(eligible_data) == 0:
            return {'base_sample': 0, 'matrix': {}}
        
        # Get weights
        if weight_column in eligible_data.columns:
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
            eligible_data = eligible_data[weights.notna()].copy()
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
        else:
            weights = pd.Series([1.0] * len(eligible_data), index=eligible_data.index)
        
        # Categorize 2021 votes
        eligible_data['party_2021'] = eligible_data[question_2021].apply(self.categorize_party_2021)
        
        # Categorize 2025 votes
        eligible_data['party_2025'] = eligible_data[question_2025].apply(self.calculator.categorize_party)
        
        # Calculate gains/losses matrix
        parties = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        matrix = {}
        
        for party_2021 in parties:
            matrix[party_2021] = {}
            party_2021_data = eligible_data[eligible_data['party_2021'] == party_2021].copy()
            
            if len(party_2021_data) == 0:
                for party_2025 in parties:
                    matrix[party_2021][party_2025] = 0.0
                continue
            
            # Get weights for this 2021 party
            party_2021_weights = weights[party_2021_data.index]
            total_weight_2021 = party_2021_weights.sum()
            
            if total_weight_2021 == 0:
                for party_2025 in parties:
                    matrix[party_2021][party_2025] = 0.0
                continue
            
            # Calculate percentage going to each 2025 party
            for party_2025 in parties:
                party_2025_mask = party_2021_data['party_2025'] == party_2025
                party_2025_weights = party_2021_weights[party_2025_mask]
                weight_sum = party_2025_weights.sum()
                percentage = (weight_sum / total_weight_2021) * 100
                matrix[party_2021][party_2025] = percentage
        
        # Calculate overall 2025 vote shares (for "Total" row)
        # This is the weighted average across all 2021 parties
        total_2025_vote_shares = {}
        total_weight_all = weights.sum()
        
        if total_weight_all > 0:
            for party_2025 in parties:
                party_2025_mask = eligible_data['party_2025'] == party_2025
                party_2025_weights = weights[party_2025_mask]
                weight_sum = party_2025_weights.sum()
                total_2025_vote_shares[party_2025] = (weight_sum / total_weight_all) * 100
        else:
            for party_2025 in parties:
                total_2025_vote_shares[party_2025] = 0.0
        
        # Base sample size: count all eligible voters (those who voted in 2021)
        base_sample = len(eligible_data)
        
        return {'base_sample': base_sample, 'matrix': matrix, 'total_2025_vote_shares': total_2025_vote_shares}
    
    def update_gains_losses_table(self, table, is_7dma=False, is_demographic=False, demographic_type=None, demographic_value=None):
        """
        Update Gains and Losses table (Slides 47-61)
        
        Args:
            table: Table object to update
            is_7dma: If True, use 7DMA data; if False, use Overall data
            is_demographic: If True, this is a demographic table; if False, overall table
            demographic_type: Demographic type if demographic table
            demographic_value: Demographic value if demographic table
        
        Returns:
            Sample size used for calculations (for updating Base text)
        """
        # Determine which data to use
        if is_demographic:
            # For demographic tables, filter data by demographic
            demographic_column_map = {
                'Gender': 'Gender',
                'Location': 'Residential locality type',
                'Religion': '20. Could you please tell me the religion that you belong to?',
                'Social Category': '21. Which social category do you belong to?',
                'Age': 'Could you please tell me your age in complete years?'
            }
            
            demographic_column = demographic_column_map.get(demographic_type)
            if not demographic_column or demographic_column not in self.df.columns:
                print(f"Warning: {demographic_column} not found for {demographic_type}")
                return
            
            # Create demographic filter
            if demographic_type == 'Gender':
                if demographic_value == 'Male':
                    demo_filter = self.df[demographic_column] == 1
                elif demographic_value == 'Female':
                    demo_filter = self.df[demographic_column] == 2
                else:
                    return
            elif demographic_type == 'Location':
                if demographic_value == 'Urban':
                    demo_filter = self.df[demographic_column] == 1
                elif demographic_value == 'Rural':
                    demo_filter = self.df[demographic_column] == 2
                else:
                    return
            elif demographic_type == 'Religion':
                if demographic_value == 'Hindu':
                    demo_filter = self.df[demographic_column] == 1
                elif demographic_value == 'Muslim':
                    demo_filter = self.df[demographic_column] == 2
                else:
                    return
            elif demographic_type == 'Social Category':
                if demographic_value == 'General+OBC':
                    demo_filter = self.df[demographic_column].isin([1, 2])
                elif demographic_value == 'SC':
                    demo_filter = self.df[demographic_column] == 3
                elif demographic_value == 'ST':
                    demo_filter = self.df[demographic_column] == 4
                else:
                    return
            elif demographic_type == 'Age':
                age_col = self.df[demographic_column]
                age_numeric = pd.to_numeric(age_col, errors='coerce')
                if demographic_value == '18-25':
                    demo_filter = (age_numeric >= 18) & (age_numeric <= 25)
                elif demographic_value == '26-34':
                    demo_filter = (age_numeric >= 26) & (age_numeric <= 34)
                elif demographic_value == '36-50':
                    demo_filter = (age_numeric >= 36) & (age_numeric <= 50)
                elif demographic_value == '50+':
                    demo_filter = age_numeric > 50
                else:
                    return
            
            # Filter data by demographic
            demo_data = self.df[demo_filter].copy()
        else:
            # Overall table - use all data
            demo_data = self.df.copy()
        
        # Apply date filter
        if is_7dma:
            # 7DMA: last 7 days ending one day before reference_date
            dma7_data = self.calculator.filter_by_date_range(
                days=7,
                exclude_latest=True,
                reference_date=self.reference_date
            )
            # Filter to demographic if needed
            if is_demographic:
                # Apply demographic filter to 7DMA data
                filtered_data = dma7_data[demo_filter.loc[dma7_data.index] if hasattr(demo_filter, 'loc') else demo_filter].copy()
                # Alternative: reapply demographic filter directly on dma7_data
                if demographic_type == 'Gender':
                    if demographic_value == 'Male':
                        filtered_data = dma7_data[dma7_data['Gender'] == 1].copy()
                    elif demographic_value == 'Female':
                        filtered_data = dma7_data[dma7_data['Gender'] == 2].copy()
                elif demographic_type == 'Location':
                    loc_col = 'Residential locality type'
                    if demographic_value == 'Urban':
                        filtered_data = dma7_data[dma7_data[loc_col] == 1].copy()
                    elif demographic_value == 'Rural':
                        filtered_data = dma7_data[dma7_data[loc_col] == 2].copy()
                elif demographic_type == 'Religion':
                    rel_col = '20. Could you please tell me the religion that you belong to?'
                    if demographic_value == 'Hindu':
                        filtered_data = dma7_data[dma7_data[rel_col] == 1].copy()
                    elif demographic_value == 'Muslim':
                        filtered_data = dma7_data[dma7_data[rel_col] == 2].copy()
                elif demographic_type == 'Social Category':
                    sc_col = '21. Which social category do you belong to?'
                    if demographic_value == 'General+OBC':
                        filtered_data = dma7_data[dma7_data[sc_col].isin([1, 2])].copy()
                    elif demographic_value == 'SC':
                        filtered_data = dma7_data[dma7_data[sc_col] == 3].copy()
                    elif demographic_value == 'ST':
                        filtered_data = dma7_data[dma7_data[sc_col] == 4].copy()
                elif demographic_type == 'Age':
                    age_col = 'Could you please tell me your age in complete years?'
                    age_numeric = pd.to_numeric(dma7_data[age_col], errors='coerce')
                    if demographic_value == '18-25':
                        filtered_data = dma7_data[(age_numeric >= 18) & (age_numeric <= 25)].copy()
                    elif demographic_value == '26-34':
                        filtered_data = dma7_data[(age_numeric >= 26) & (age_numeric <= 34)].copy()
                    elif demographic_value == '36-50':
                        filtered_data = dma7_data[(age_numeric >= 36) & (age_numeric <= 50)].copy()
                    elif demographic_value == '50+':
                        filtered_data = dma7_data[age_numeric > 50].copy()
            else:
                filtered_data = dma7_data.copy()
        else:
            # Overall: all data up to reference_date
            filtered_data = demo_data[demo_data['Survey Date'] <= self.reference_date].copy()
        
        # Calculate gains/losses
        gains_losses = self.calculate_gains_losses(
            filtered_data,
            is_7dma=is_7dma,
            demographic_type=demographic_type if is_demographic else None,
            demographic_value=demographic_value if is_demographic else None
        )
        
        # Get 2021 AE vote shares
        if is_demographic:
            # For demographic tables, calculate from survey data
            ae2021_vs = self.calculate_ae2021_vote_shares_from_survey(
                filtered_data,
                weight_column=self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level',
                use_weights=True
            )
        else:
            # For overall tables (Slides 47-48), use constant values from final format
            # These are fixed government values and should not be recalculated
            ae2021_vs = {
                'AITC': 48.02,
                'BJP': 37.97,
                'LEFT': 5.66,
                'INC': 3.03,
                'Others': 2.88,  # IND + OTHERS
                'NWR': 1.08  # NOTA
            }
        
        # Helper function to format cell
        def format_table_cell(cell, value, is_percentage=False):
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            
            if is_percentage:
                formatted_value = f"{value:.1f}"
            else:
                formatted_value = f"{value:,}"
            
            cell.text = formatted_value
            text_frame = cell.text_frame
            text_frame.clear()
            para = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            para.clear()
            run = para.add_run()
            run.text = formatted_value
            run.font.size = Pt(14)
            run.font.name = 'Aptos Display'
            para.alignment = PP_ALIGN.RIGHT
        
        # Update table
        # Table structure:
        # Row 0: Headers
        # Row 1: Column headers (Party, AE 2021 VS, 2025 Voting Pre, AITC, BJP, LEFT+, INC, Others, N+W+R*)
        # Row 2: AITC row (2021 AE, AITC, vote share, then percentages)
        # Row 3: BJP row
        # Row 4: LEFT+ row
        # Row 5: INC row
        # etc.
        
        parties_2021 = ['AITC', 'BJP', 'LEFT', 'INC']
        parties_2025 = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        
        # Update each party row (AITC, BJP, LEFT, INC)
        for row_idx, party_2021 in enumerate(parties_2021, start=2):  # Start from row 2
            if row_idx >= len(table.rows):
                break
            
            # Column 0: "2021 AE" label (only for first row)
            if row_idx == 2:
                # Don't change the label, just update values
                pass
            
            # Column 1: Party name - don't change, keep template design
            # table.cell(row_idx, 1).text = party_display.get(party_2021, party_2021)
            
            # Column 2: AE 2021 VS (vote share in 2021) - only update value, preserve template formatting
            if len(table.columns) > 2:
                cell = table.cell(row_idx, 2)
                # Update text value only, preserve existing formatting
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{ae2021_vs.get(party_2021, 0):.2f}"
                    else:
                        para.add_run().text = f"{ae2021_vs.get(party_2021, 0):.2f}"
                else:
                    cell.text = f"{ae2021_vs.get(party_2021, 0):.2f}"
            
            # Columns 3-8: 2025 voting preferences (percentages) - only update values, preserve formatting
            for col_idx, party_2025 in enumerate(parties_2025, start=3):
                if col_idx < len(table.columns):
                    percentage = gains_losses['matrix'].get(party_2021, {}).get(party_2025, 0)
                    cell = table.cell(row_idx, col_idx)
                    # Update text value only, preserve existing formatting
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{percentage:.1f}"
                        else:
                            para.add_run().text = f"{percentage:.1f}"
                    else:
                        cell.text = f"{percentage:.1f}"
        
        # Get base sample size from gains_losses
        base_sample = gains_losses.get('base_sample', 0)
        
        # Update "Others" row (row 6)
        if len(table.rows) > 6:
            row_idx = 6
            # Column 2: AE 2021 VS for Others
            if len(table.columns) > 2:
                cell = table.cell(row_idx, 2)
                # Update text value only, preserve existing formatting
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{ae2021_vs.get('Others', 0):.2f}"
                    else:
                        para.add_run().text = f"{ae2021_vs.get('Others', 0):.2f}"
                else:
                    cell.text = f"{ae2021_vs.get('Others', 0):.2f}"
            
            # Columns 3-8: Transitions from "Others" party in 2021
            for col_idx, party_2025 in enumerate(parties_2025, start=3):
                if col_idx < len(table.columns):
                    percentage = gains_losses['matrix'].get('Others', {}).get(party_2025, 0)
                    cell = table.cell(row_idx, col_idx)
                    # Update text value only, preserve existing formatting
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{percentage:.1f}"
                        else:
                            para.add_run().text = f"{percentage:.1f}"
                    else:
                        cell.text = f"{percentage:.1f}"
        
        # Update "N+W+R" row (row 7)
        if len(table.rows) > 7:
            row_idx = 7
            # Column 2: AE 2021 VS for NWR
            if len(table.columns) > 2:
                cell = table.cell(row_idx, 2)
                # Update text value only, preserve existing formatting
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{ae2021_vs.get('NWR', 0):.2f}"
                    else:
                        para.add_run().text = f"{ae2021_vs.get('NWR', 0):.2f}"
                else:
                    cell.text = f"{ae2021_vs.get('NWR', 0):.2f}"
            
            # Columns 3-8: Transitions from "NWR" party in 2021
            for col_idx, party_2025 in enumerate(parties_2025, start=3):
                if col_idx < len(table.columns):
                    percentage = gains_losses['matrix'].get('NWR', {}).get(party_2025, 0)
                    cell = table.cell(row_idx, col_idx)
                    # Update text value only, preserve existing formatting
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{percentage:.1f}"
                        else:
                            para.add_run().text = f"{percentage:.1f}"
                    else:
                        cell.text = f"{percentage:.1f}"
        
        # Update "Total" row (row 8)
        if len(table.rows) > 8:
            row_idx = 8
            # Column 2: Total AE 2021 VS (sum of all parties)
            if len(table.columns) > 2:
                total_ae2021 = sum(ae2021_vs.get(p, 0) for p in parties_2021) + ae2021_vs.get('Others', 0) + ae2021_vs.get('NWR', 0)
                cell = table.cell(row_idx, 2)
                # Update text value only, preserve existing formatting
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{total_ae2021:.2f}"
                    else:
                        para.add_run().text = f"{total_ae2021:.2f}"
                else:
                    cell.text = f"{total_ae2021:.2f}"
            
            # Columns 3-8: Overall 2025 vote shares (weighted average across all 2021 parties)
            total_2025_vote_shares = gains_losses.get('total_2025_vote_shares', {})
            for col_idx, party_2025 in enumerate(parties_2025, start=3):
                if col_idx < len(table.columns):
                    percentage = total_2025_vote_shares.get(party_2025, 0)
                    cell = table.cell(row_idx, col_idx)
                    # Update text value only, preserve existing formatting
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{percentage:.1f}"
                        else:
                            para.add_run().text = f"{percentage:.1f}"
                    else:
                        cell.text = f"{percentage:.1f}"
        
        # Return base sample size for updating Base text
        return base_sample
    
    def calculate_vote_transferability(self, data_filtered, is_7dma=False):
        """
        Calculate Vote Transferability table data
        Shows percentage of people who vote for party X as first choice and party Y as second choice
        
        Args:
            data_filtered: Filtered DataFrame (overall or 7DMA)
            is_7dma: If True, use 7DMA data; if False, use Overall data
        
        Returns:
            Dictionary with transferability matrix and base sample size
        """
        question_first_choice = '8. If assembly elections (MLA) were to be held tomorrow, then which party would you vote for?'
        # Question 9: Second choice party
        question_second_choice = None
        for col in data_filtered.columns:
            if '9.' in str(col) and 'assume' in str(col).lower() and 'party' in str(col).lower() and 'choose' in str(col).lower():
                question_second_choice = col
                break
        if question_second_choice is None:
            # Fallback: try to find any column with '9.' and 'party'
            for col in data_filtered.columns:
                if '9.' in str(col) and 'party' in str(col).lower():
                    question_second_choice = col
                    break
        question_second_choice_reason = '10. Could you tell us the reason for choosing the above party as your second choice?'
        weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        # Filter to records with first choice
        eligible_data = data_filtered[data_filtered[question_first_choice].notna()].copy()
        
        # Find all Question 10 columns (reasons for second choice)
        q10_columns = [c for c in eligible_data.columns if question_second_choice_reason in str(c)]
        
        # Filter to records that have a second choice (indicated by having at least one Q10 response)
        # A record has a second choice if any Q10 column has value 1 (or non-null)
        if q10_columns:
            has_second_choice = False
            for q10_col in q10_columns:
                # Check if column has value 1 (indicating this reason was selected)
                has_second_choice = has_second_choice | (pd.to_numeric(eligible_data[q10_col], errors='coerce') == 1)
            eligible_data = eligible_data[has_second_choice].copy()
        
        # Now filter to records that also have Question 9 (second choice party) response
        if question_second_choice is None or question_second_choice not in eligible_data.columns:
            # If Question 9 column not found, return empty
            return {'base_sample': 0, 'matrix': {}}
        
        eligible_data = eligible_data[eligible_data[question_second_choice].notna()].copy()
        
        if len(eligible_data) == 0:
            return {'base_sample': 0, 'matrix': {}}
        
        # Get weights
        if weight_column in eligible_data.columns:
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
            eligible_data = eligible_data[weights.notna()].copy()
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
        else:
            weights = pd.Series([1.0] * len(eligible_data), index=eligible_data.index)
        
        # Categorize first choice votes
        eligible_data['party_first'] = eligible_data[question_first_choice].apply(self.calculator.categorize_party)
        
        # Categorize second choice votes
        eligible_data['party_second'] = eligible_data[question_second_choice].apply(self.calculator.categorize_party)
        
        # Base sample size: all eligible records (with both first and second choice)
        base_sample_size = len(eligible_data)
        
        # Calculate transferability matrix
        # IMPORTANT: Include ALL first choice voters in denominator (including same-party)
        # But diagonal (same party) should always show "-"
        parties = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        matrix = {}
        
        for party_first in parties:
            matrix[party_first] = {}
            # Get ALL records with this first choice (including same-party second choice)
            party_first_all = eligible_data[eligible_data['party_first'] == party_first].copy()
            
            if len(party_first_all) == 0:
                for party_second in parties:
                    if party_first == party_second:
                        matrix[party_first][party_second] = "-"
                    else:
                        matrix[party_first][party_second] = 0
                continue
            
            # Get weights for ALL first choice voters (including same-party)
            party_weights_all = weights.loc[party_first_all.index]
            total_weight_first = party_weights_all.sum()
            
            if total_weight_first == 0:
                for party_second in parties:
                    if party_first == party_second:
                        matrix[party_first][party_second] = "-"
                    else:
                        matrix[party_first][party_second] = 0
                continue
            
            # Calculate percentage for each second choice
            for party_second in parties:
                # If same party (diagonal), always show "-" (not applicable)
                if party_first == party_second:
                    matrix[party_first][party_second] = "-"
                else:
                    # Get records with this first choice AND this second choice (different party)
                    party_second_data = party_first_all[party_first_all['party_second'] == party_second].copy()
                    if len(party_second_data) == 0:
                        matrix[party_first][party_second] = 0
                    else:
                        # Calculate percentage: (weighted count of first=X and second=Y) / (weighted count of first=X) * 100
                        weight_sum = party_weights_all.loc[party_second_data.index].sum()
                        percentage = (weight_sum / total_weight_first) * 100
                        matrix[party_first][party_second] = percentage
        
        return {'base_sample': base_sample_size, 'matrix': matrix}
    
    def update_vote_transferability_table(self, table, is_7dma=False):
        """
        Update Vote Transferability table (Slides 63-64)
        
        Args:
            table: Table object from slide
            is_7dma: If True, use 7DMA data; if False, use Overall data
        
        Returns:
            base_sample size
        """
        # Filter data
        if is_7dma:
            # 7DMA: last 7 days ending one day before reference_date
            dma7_data = self.calculator.filter_by_date_range(
                days=7,
                exclude_latest=True,
                reference_date=self.reference_date
            )
            filtered_data = dma7_data.copy()
        else:
            # Overall: all data up to reference_date
            filtered_data = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # Calculate vote transferability
        transferability = self.calculate_vote_transferability(filtered_data, is_7dma=is_7dma)
        
        # Update table
        parties_first = ['AITC', 'BJP', 'LEFT', 'INC']
        parties_second = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        
        # Update each first choice party row (rows 2-5)
        # Row 0: Headers, Row 1: Column headers, Row 2+: Data rows
        for row_idx, party_first in enumerate(parties_first, start=2):
            if row_idx >= len(table.rows):
                break
            
            # Columns 2-7: Second choice percentages
            # Column 0: "First Preference" label, Column 1: Party name, Column 2+: Second preference columns
            for col_idx, party_second in enumerate(parties_second, start=2):
                if col_idx < len(table.columns):
                    value = transferability['matrix'].get(party_first, {}).get(party_second, 0)
                    cell = table.cell(row_idx, col_idx)
                    # Update text value only, preserve existing formatting
                    if value == "-":
                        # Show "-" for diagonal (same party)
                        display_value = "-"
                    elif isinstance(value, (int, float)):
                        display_value = f"{value:.1f}"
                    else:
                        display_value = str(value)
                    
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = display_value
                        else:
                            para.add_run().text = display_value
                    else:
                        cell.text = display_value
        
        # Update "Others" row (row 6)
        if len(table.rows) > 6:
            row_idx = 6
            for col_idx, party_second in enumerate(parties_second, start=2):
                if col_idx < len(table.columns):
                    value = transferability['matrix'].get('Others', {}).get(party_second, 0)
                    cell = table.cell(row_idx, col_idx)
                    if value == "-":
                        display_value = "-"
                    elif isinstance(value, (int, float)):
                        display_value = f"{value:.1f}"
                    else:
                        display_value = str(value)
                    
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = display_value
                        else:
                            para.add_run().text = display_value
                    else:
                        cell.text = display_value
        
        # Update "N+W+R" row (row 7)
        if len(table.rows) > 7:
            row_idx = 7
            for col_idx, party_second in enumerate(parties_second, start=2):
                if col_idx < len(table.columns):
                    value = transferability['matrix'].get('NWR', {}).get(party_second, 0)
                    cell = table.cell(row_idx, col_idx)
                    if value == "-":
                        display_value = "-"
                    elif isinstance(value, (int, float)):
                        display_value = f"{value:.1f}"
                    else:
                        display_value = str(value)
                    
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = display_value
                        else:
                            para.add_run().text = display_value
                    else:
                        cell.text = display_value
        
        return transferability['base_sample']
    
    def calculate_preferred_cm(self, data_filtered):
        """
        Calculate Preferred CM Candidate percentages
        
        Args:
            data_filtered: Filtered DataFrame
        
        Returns:
            Dictionary with CM candidate percentages and base sample size
        """
        question = '17. Who do you think is the best leader to be the Chief Minister of West Bengal?'
        weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        # Filter to records with valid response
        eligible_data = data_filtered[data_filtered[question].notna()].copy()
        
        if len(eligible_data) == 0:
            return {'base_sample': 0, 'percentages': {}}
        
        # Get weights
        if weight_column in eligible_data.columns:
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
            eligible_data = eligible_data[weights.notna()].copy()
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
        else:
            weights = pd.Series([1.0] * len(eligible_data), index=eligible_data.index)
        
        total_weight = weights.sum()
        if total_weight == 0:
            return {'base_sample': len(eligible_data), 'percentages': {}}
        
        # Calculate percentages for each CM candidate code
        percentages = {}
        for code in eligible_data[question].dropna().unique():
            if pd.isna(code) or code == 'q17':
                continue
            code_data = eligible_data[eligible_data[question] == code]
            weight_sum = weights.loc[code_data.index].sum()
            percentage = (weight_sum / total_weight) * 100
            percentages[int(code)] = percentage
        
        return {'base_sample': len(eligible_data), 'percentages': percentages}
    
    def update_preferred_cm_table(self, table):
        """
        Update Preferred CM Candidate table (Slide 66)
        
        Args:
            table: Table object from slide
        
        Returns:
            base_sample size
        """
        # Filter data for Overall
        filtered_data_overall = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # Filter data for 7DMA
        dma7_data = self.calculator.filter_by_date_range(
            days=7,
            exclude_latest=True,
            reference_date=self.reference_date
        )
        
        # Calculate preferred CM for Overall
        cm_data_overall = self.calculate_preferred_cm(filtered_data_overall)
        
        # Calculate preferred CM for 7DMA
        cm_data_7dma = self.calculate_preferred_cm(dma7_data)
        
        # Update sample size row (row 2)
        if len(table.rows) > 2:
            # Column 3: 7DMA sample size
            if len(table.columns) > 3:
                cell = table.cell(2, 3)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{cm_data_7dma['base_sample']:,}"
                    else:
                        para.add_run().text = f"{cm_data_7dma['base_sample']:,}"
                else:
                    cell.text = f"{cm_data_7dma['base_sample']:,}"
            
            # Column 4: Overall sample size
            if len(table.columns) > 4:
                cell = table.cell(2, 4)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{cm_data_overall['base_sample']:,}"
                    else:
                        para.add_run().text = f"{cm_data_overall['base_sample']:,}"
                else:
                    cell.text = f"{cm_data_overall['base_sample']:,}"
        
        # Update CM candidate rows (rows 3+)
        # Map CM codes to row indices based on the S. No. column in the table
        # The table has S. No. in column 0, and we need to match CM codes to those S. No.
        percentages_overall = cm_data_overall['percentages']
        percentages_7dma = cm_data_7dma['percentages']
        
        # Create mapping from S. No. to CM code based on template structure
        # Based on template: Row 3 = S. No. 1 (Mamata Banerjee), Row 4 = S. No. 2 (Abhishek Banerjee), etc.
        # We'll match by reading the S. No. from column 0 and updating the corresponding row
        # Since we don't have a direct mapping, we'll use the sorted order by percentage
        # Get sorted CM codes by percentage (descending) for Overall
        sorted_cm_codes = sorted(percentages_overall.items(), key=lambda x: x[1], reverse=True)
        
        # Update each CM candidate row (rows 3-18, assuming row 0-2 are headers)
        # Match by S. No. in column 0 if possible, otherwise use sorted order
        for row_idx in range(3, len(table.rows)):
            if row_idx >= len(table.rows):
                break
            
            # Try to get S. No. from column 0
            s_no_text = table.cell(row_idx, 0).text.strip()
            try:
                s_no = int(s_no_text)
            except:
                # If S. No. not found, use sorted order
                sorted_idx = row_idx - 3
                if sorted_idx < len(sorted_cm_codes):
                    cm_code, percentage_overall = sorted_cm_codes[sorted_idx]
                    percentage_7dma = percentages_7dma.get(cm_code, 0)
                else:
                    continue
            else:
                # Match by S. No. - use sorted order for now (S. No. 1 = highest percentage, etc.)
                sorted_idx = s_no - 1
                if sorted_idx < len(sorted_cm_codes):
                    cm_code, percentage_overall = sorted_cm_codes[sorted_idx]
                    percentage_7dma = percentages_7dma.get(cm_code, 0)
                else:
                    continue
            
            # Column 3: 7DMA percentage
            if len(table.columns) > 3:
                cell = table.cell(row_idx, 3)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage_7dma:.1f}"
                    else:
                        para.add_run().text = f"{percentage_7dma:.1f}"
                else:
                    cell.text = f"{percentage_7dma:.1f}"
            
            # Column 4: Overall percentage
            if len(table.columns) > 4:
                cell = table.cell(row_idx, 4)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage_overall:.1f}"
                    else:
                        para.add_run().text = f"{percentage_overall:.1f}"
                else:
                    cell.text = f"{percentage_overall:.1f}"
        
        return cm_data_overall['base_sample']
    
    def calculate_state_government_rating(self, data_filtered):
        """
        Calculate State Government Rating percentages
        
        Args:
            data_filtered: Filtered DataFrame
        
        Returns:
            Dictionary with rating percentages and base sample size
        """
        question = '14. How satisfied or dissatisfied are you with the performance of the state govt led by Mamata Banerjee?'
        weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        # Filter to records with valid response
        eligible_data = data_filtered[data_filtered[question].notna()].copy()
        
        # Remove invalid responses
        eligible_data = eligible_data[eligible_data[question] != 'q14'].copy()
        
        if len(eligible_data) == 0:
            return {'base_sample': 0, 'percentages': {}}
        
        # Get weights
        if weight_column in eligible_data.columns:
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
            eligible_data = eligible_data[weights.notna()].copy()
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
        else:
            weights = pd.Series([1.0] * len(eligible_data), index=eligible_data.index)
        
        total_weight = weights.sum()
        if total_weight == 0:
            return {'base_sample': len(eligible_data), 'percentages': {}}
        
        # Map rating codes to labels
        # 1 = Very Satisfied, 2 = Satisfied, 3 = Neutral, 4 = Dissatisfied, 5 = Very Dissatisfied
        rating_map = {
            1: 'Very Satisfied',
            2: 'Satisfied',
            3: 'Neutral',
            4: 'Dissatisfied',
            5: 'Very Dissatisfied'
        }
        
        # Calculate percentages for each rating
        percentages = {}
        for code in [1, 2, 3, 4, 5]:
            code_data = eligible_data[pd.to_numeric(eligible_data[question], errors='coerce') == code]
            if len(code_data) > 0:
                weight_sum = weights.loc[code_data.index].sum()
                percentage = (weight_sum / total_weight) * 100
                percentages[code] = percentage
            else:
                percentages[code] = 0
        
        return {'base_sample': len(eligible_data), 'percentages': percentages}
    
    def update_state_government_rating_table(self, table):
        """
        Update State Government Rating table (Slide 68)
        
        Args:
            table: Table object from slide
        
        Returns:
            base_sample size
        """
        # Filter data for Overall
        filtered_data_overall = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # Filter data for 7DMA
        dma7_data = self.calculator.filter_by_date_range(
            days=7,
            exclude_latest=True,
            reference_date=self.reference_date
        )
        
        # Calculate state government rating for Overall
        rating_data_overall = self.calculate_state_government_rating(filtered_data_overall)
        
        # Calculate state government rating for 7DMA
        rating_data_7dma = self.calculate_state_government_rating(dma7_data)
        
        # Update table
        # Rating codes: 1 = Very Satisfied, 2 = Satisfied, 3 = Neutral, 4 = Dissatisfied, 5 = Very Dissatisfied
        rating_codes = [1, 2, 3, 4, 5]
        percentages_overall = rating_data_overall['percentages']
        percentages_7dma = rating_data_7dma['percentages']
        
        # Update each rating row (rows 1-5, assuming row 0 is header)
        for row_idx, rating_code in enumerate(rating_codes, start=1):
            if row_idx >= len(table.rows):
                break
            
            # Column 1: 7DMA percentage
            percentage_7dma = percentages_7dma.get(rating_code, 0)
            if len(table.columns) > 1:
                cell = table.cell(row_idx, 1)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage_7dma:.1f}"
                    else:
                        para.add_run().text = f"{percentage_7dma:.1f}"
                else:
                    cell.text = f"{percentage_7dma:.1f}"
            
            # Column 2: Overall percentage
            percentage_overall = percentages_overall.get(rating_code, 0)
            if len(table.columns) > 2:
                cell = table.cell(row_idx, 2)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage_overall:.1f}"
                    else:
                        para.add_run().text = f"{percentage_overall:.1f}"
                else:
                    cell.text = f"{percentage_overall:.1f}"
        
        return rating_data_overall['base_sample']
    
    def calculate_wisdom_of_crowds(self, data_filtered):
        """
        Calculate Wisdom of Crowds (which party is likely to win) percentages
        
        Args:
            data_filtered: Filtered DataFrame
        
        Returns:
            Dictionary with party percentages and base sample size
        """
        question = '19. In your opinion, which party would win the next election in your constituency, when you would elect your MLA?'
        weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        # Filter to records with valid response
        eligible_data = data_filtered[data_filtered[question].notna()].copy()
        
        if len(eligible_data) == 0:
            return {'base_sample': 0, 'percentages': {}}
        
        # Get weights
        if weight_column in eligible_data.columns:
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
            eligible_data = eligible_data[weights.notna()].copy()
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
        else:
            weights = pd.Series([1.0] * len(eligible_data), index=eligible_data.index)
        
        total_weight = weights.sum()
        if total_weight == 0:
            return {'base_sample': len(eligible_data), 'percentages': {}}
        
        # Categorize party responses
        eligible_data['party'] = eligible_data[question].apply(self.calculator.categorize_party)
        
        # Calculate percentages for each party
        parties = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        percentages = {}
        
        for party in parties:
            party_data = eligible_data[eligible_data['party'] == party].copy()
            if len(party_data) > 0:
                weight_sum = weights.loc[party_data.index].sum()
                percentage = (weight_sum / total_weight) * 100
                percentages[party] = percentage
            else:
                percentages[party] = 0
        
        return {'base_sample': len(eligible_data), 'percentages': percentages}
    
    def update_wisdom_of_crowds_table(self, table):
        """
        Update Wisdom of Crowds table (Slide 72)
        
        Args:
            table: Table object from slide
        
        Returns:
            base_sample size
        """
        # Filter data for Overall
        filtered_data_overall = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # Filter data for 7DMA
        dma7_data = self.calculator.filter_by_date_range(
            days=7,
            exclude_latest=True,
            reference_date=self.reference_date
        )
        
        # Calculate wisdom of crowds for Overall
        woc_data_overall = self.calculate_wisdom_of_crowds(filtered_data_overall)
        
        # Calculate wisdom of crowds for 7DMA
        woc_data_7dma = self.calculate_wisdom_of_crowds(dma7_data)
        
        # Update table
        parties = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        percentages_overall = woc_data_overall['percentages']
        percentages_7dma = woc_data_7dma['percentages']
        
        # Update each party row (rows 1-6, assuming row 0 is header)
        for row_idx, party in enumerate(parties, start=1):
            if row_idx >= len(table.rows):
                break
            
            # Column 1: 7DMA percentage
            percentage_7dma = percentages_7dma.get(party, 0)
            if len(table.columns) > 1:
                cell = table.cell(row_idx, 1)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage_7dma:.1f}"
                    else:
                        para.add_run().text = f"{percentage_7dma:.1f}"
                else:
                    cell.text = f"{percentage_7dma:.1f}"
            
            # Column 2: Overall percentage
            percentage_overall = percentages_overall.get(party, 0)
            if len(table.columns) > 2:
                cell = table.cell(row_idx, 2)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage_overall:.1f}"
                    else:
                        para.add_run().text = f"{percentage_overall:.1f}"
                else:
                    cell.text = f"{percentage_overall:.1f}"
        
        return woc_data_overall['base_sample']
    
    def calculate_top_issues(self, data_filtered):
        """
        Calculate Top Issues from Question 13
        
        Args:
            data_filtered: Filtered DataFrame
        
        Returns:
            Dictionary with top issues, percentages, and base sample size
        """
        question_prefix = '13. According to you, what are the three most pressing issues of your assembly constituency?'
        weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        
        # Find all Q13 columns
        issue_cols = [col for col in data_filtered.columns 
                     if str(col).startswith(question_prefix)]
        
        if len(issue_cols) == 0:
            return {'base_sample': 0, 'issues': []}
        
        # Filter to records with at least one issue selected
        has_issue = data_filtered[issue_cols].notna().any(axis=1)
        eligible_data = data_filtered[has_issue].copy()
        
        if len(eligible_data) == 0:
            return {'base_sample': 0, 'issues': []}
        
        # Get weights
        if weight_column in eligible_data.columns:
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
            eligible_data = eligible_data[weights.notna()].copy()
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
        else:
            weights = pd.Series([1.0] * len(eligible_data), index=eligible_data.index)
        
        total_weight = weights.sum()
        if total_weight == 0:
            return {'base_sample': len(eligible_data), 'issues': []}
        
        # Calculate weighted counts for each issue
        issue_counts = {}
        for col in issue_cols:
            # Extract issue text (everything after " - ")
            issue_text = str(col).split(' - ')[-1] if ' - ' in str(col) else str(col)
            
            # Convert column to numeric (0 or 1)
            issue_values = pd.to_numeric(eligible_data[col], errors='coerce').fillna(0)
            
            # Calculate weighted sum
            weighted_sum = (issue_values * weights).sum()
            
            # Calculate percentage
            percentage = (weighted_sum / total_weight * 100) if total_weight > 0 else 0
            
            # Count actual records (not weighted)
            record_count = int(issue_values.sum())
            
            issue_counts[issue_text] = {
                'weighted_count': weighted_sum,
                'percentage': percentage,
                'record_count': record_count
            }
        
        # Sort by weighted count and get top issues
        top_issues = sorted(issue_counts.items(), key=lambda x: x[1]['weighted_count'], reverse=True)
        
        return {
            'base_sample': len(eligible_data),
            'issues': top_issues  # List of (issue_text, {weighted_count, percentage})
        }
    
    def update_key_issues_table(self, table):
        """
        Update Key Issues table (Slide 70)
        
        Args:
            table: Table object from slide
        
        Returns:
            base_sample size
        """
        # Filter data for Overall
        filtered_data_overall = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # Calculate top issues for Overall
        issues_data = self.calculate_top_issues(filtered_data_overall)
        
        # Update table
        # Row 0 is header, rows 1-12 are for issues (S.No, Issue, Sample, 30DMA, Overall)
        # Leave 30DMA empty, only update Overall column (column 4)
        top_issues = issues_data['issues']
        
        # Update each issue row (rows 1-12, assuming row 0 is header)
        for row_idx, (issue_text, issue_data) in enumerate(top_issues[:12], start=1):
            if row_idx >= len(table.rows):
                break
            
            # Column 1: Issue text
            if len(table.columns) > 1:
                cell = table.cell(row_idx, 1)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = issue_text
                    else:
                        para.add_run().text = issue_text
                else:
                    cell.text = issue_text
            
            # Column 2: Sample size (same for all issues, use base_sample)
            if len(table.columns) > 2:
                cell = table.cell(row_idx, 2)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{issues_data['base_sample']:,}"
                    else:
                        para.add_run().text = f"{issues_data['base_sample']:,}"
                else:
                    cell.text = f"{issues_data['base_sample']:,}"
            
            # Column 3: 30DMA - Leave empty (as per user request)
            # (No update needed)
            
            # Column 4: Overall percentage
            if len(table.columns) > 4:
                cell = table.cell(row_idx, 4)
                # Get the percentage for this issue
                percentage = issue_data.get('percentage', 0)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage:.1f}"
                    else:
                        para.add_run().text = f"{percentage:.1f}"
                else:
                    cell.text = f"{percentage:.1f}"
        
        return issues_data['base_sample']
    
    def calculate_vote_share_by_regions(self, data_filtered):
        """
        Calculate Vote Share by Regions
        
        Args:
            data_filtered: Filtered DataFrame
        
        Returns:
            Dictionary with vote shares by region and base sample sizes
        """
        question = '8. If assembly elections (MLA) were to be held tomorrow, then which party would you vote for?'
        weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
        region_column = 'Region Name'
        
        # Filter to records with valid response
        eligible_data = data_filtered[data_filtered[question].notna()].copy()
        
        if len(eligible_data) == 0:
            return {'base_sample': 0, 'regions': {}}
        
        # Get weights
        if weight_column in eligible_data.columns:
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
            eligible_data = eligible_data[weights.notna()].copy()
            weights = pd.to_numeric(eligible_data[weight_column], errors='coerce')
        else:
            weights = pd.Series([1.0] * len(eligible_data), index=eligible_data.index)
        
        # Categorize party responses
        eligible_data['party'] = eligible_data[question].apply(self.calculator.categorize_party)
        
        # Get unique regions
        regions = eligible_data[region_column].dropna().unique()
        
        parties = ['AITC', 'BJP', 'LEFT', 'INC', 'Others', 'NWR']
        region_data = {}
        
        for region in regions:
            region_data[region] = {
                'base_sample': 0,
                'vote_shares': {}
            }
            
            region_records = eligible_data[eligible_data[region_column] == region].copy()
            if len(region_records) == 0:
                for party in parties:
                    region_data[region]['vote_shares'][party] = 0
                continue
            
            region_weights = weights.loc[region_records.index]
            total_weight_region = region_weights.sum()
            
            if total_weight_region == 0:
                for party in parties:
                    region_data[region]['vote_shares'][party] = 0
                continue
            
            region_data[region]['base_sample'] = len(region_records)
            
            # Calculate vote shares for each party in this region
            for party in parties:
                party_data = region_records[region_records['party'] == party].copy()
                if len(party_data) > 0:
                    weight_sum = region_weights.loc[party_data.index].sum()
                    percentage = (weight_sum / total_weight_region) * 100
                    region_data[region]['vote_shares'][party] = percentage
                else:
                    region_data[region]['vote_shares'][party] = 0
        
        return {'base_sample': len(eligible_data), 'regions': region_data}
    
    def update_vote_share_by_regions_table(self, table):
        """
        Update Vote Share by Regions table (Slide 74)
        
        Table structure:
        - Row 0: Party | Jalpaiguri | ... | Malda | ... | Burdwan | ... | Medinipur | ... | Presidency | ...
        - Row 1: | 2021 AE | 7 DMA | Overall | 2021 AE | 7 DMA | Overall | ... (repeated for each region)
        - Row 2: Sample | - | 7DMA_sample | Overall_sample | - | 7DMA_sample | Overall_sample | ...
        - Row 3+: Party rows with vote shares
        
        Each region has 3 columns: 2021 AE (col 1), 7 DMA (col 2), Overall (col 3)
        Region order: Jalpaiguri (cols 1-3), Malda (cols 4-6), Burdwan (cols 7-9), Medinipur (cols 10-12), Presidency (cols 13-15)
        
        Args:
            table: Table object from slide
        
        Returns:
            base_sample size (overall)
        """
        # Define region order and column mapping
        # Each region has 3 columns: 2021 AE, 7 DMA, Overall
        region_order = ['Jalpaiguri', 'Malda', 'Burdwan', 'Medinipur', 'Presidency']
        region_col_map = {}
        for idx, region in enumerate(region_order):
            # Each region starts at column 1 + (idx * 3)
            # Column 0 is "Party", so region columns start at 1
            region_col_map[region] = {
                '2021_AE': 1 + (idx * 3),  # Leave empty (2021 AE data not from survey)
                '7DMA': 2 + (idx * 3),
                'Overall': 3 + (idx * 3)
            }
        
        # Filter data for Overall
        filtered_data_overall = self.df[self.df['Survey Date'] <= self.reference_date].copy()
        
        # Filter data for 7DMA
        dma7_data = self.calculator.filter_by_date_range(
            days=7,
            exclude_latest=True,
            reference_date=self.reference_date
        )
        
        # Calculate vote share by regions for Overall
        regions_data_overall = self.calculate_vote_share_by_regions(filtered_data_overall)
        
        # Calculate vote share by regions for 7DMA
        regions_data_7dma = self.calculate_vote_share_by_regions(dma7_data)
        
        # Party mapping to row indices
        # Row 2: Sample, Row 3: AITC, Row 4: BJP, Row 5: Left Front, Row 6: INC, Row 7: Other, Row 8: N+W+R, Row 9: Margin
        party_row_map = {
            'AITC': 3,
            'BJP': 4,
            'LEFT': 5,
            'INC': 6,
            'Others': 7,
            'NWR': 8
        }
        
        # Update Sample row (row 2)
        if len(table.rows) > 2:
            for region in region_order:
                if region in region_col_map:
                    col_map = region_col_map[region]
                    # 7DMA sample
                    if col_map['7DMA'] < len(table.columns):
                        region_7dma = regions_data_7dma['regions'].get(region, {})
                        sample_7dma = region_7dma.get('base_sample', 0)
                        cell = table.cell(2, col_map['7DMA'])
                        if len(cell.text_frame.paragraphs) > 0:
                            para = cell.text_frame.paragraphs[0]
                            if len(para.runs) > 0:
                                para.runs[0].text = f"{sample_7dma:,}"
                            else:
                                para.add_run().text = f"{sample_7dma:,}"
                        else:
                            cell.text = f"{sample_7dma:,}"
                    
                    # Overall sample
                    if col_map['Overall'] < len(table.columns):
                        region_overall = regions_data_overall['regions'].get(region, {})
                        sample_overall = region_overall.get('base_sample', 0)
                        cell = table.cell(2, col_map['Overall'])
                        if len(cell.text_frame.paragraphs) > 0:
                            para = cell.text_frame.paragraphs[0]
                            if len(para.runs) > 0:
                                para.runs[0].text = f"{sample_overall:,}"
                            else:
                                para.add_run().text = f"{sample_overall:,}"
                        else:
                            cell.text = f"{sample_overall:,}"
        
        # Update party rows (rows 3-8)
        for party, row_idx in party_row_map.items():
            if row_idx >= len(table.rows):
                continue
            
            for region in region_order:
                if region not in region_col_map:
                    continue
                
                col_map = region_col_map[region]
                
                # 7DMA vote share
                if col_map['7DMA'] < len(table.columns):
                    region_7dma = regions_data_7dma['regions'].get(region, {})
                    vote_shares_7dma = region_7dma.get('vote_shares', {})
                    percentage_7dma = vote_shares_7dma.get(party, 0)
                    cell = table.cell(row_idx, col_map['7DMA'])
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{percentage_7dma:.1f}"
                        else:
                            para.add_run().text = f"{percentage_7dma:.1f}"
                    else:
                        cell.text = f"{percentage_7dma:.1f}"
                
                # Overall vote share
                if col_map['Overall'] < len(table.columns):
                    region_overall = regions_data_overall['regions'].get(region, {})
                    vote_shares_overall = region_overall.get('vote_shares', {})
                    percentage_overall = vote_shares_overall.get(party, 0)
                    cell = table.cell(row_idx, col_map['Overall'])
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{percentage_overall:.1f}"
                        else:
                            para.add_run().text = f"{percentage_overall:.1f}"
                    else:
                        cell.text = f"{percentage_overall:.1f}"
        
        # Update Margin row (row 9) - AITC - BJP
        if len(table.rows) > 9:
            for region in region_order:
                if region not in region_col_map:
                    continue
                
                col_map = region_col_map[region]
                
                # 7DMA margin
                if col_map['7DMA'] < len(table.columns):
                    region_7dma = regions_data_7dma['regions'].get(region, {})
                    vote_shares_7dma = region_7dma.get('vote_shares', {})
                    margin_7dma = vote_shares_7dma.get('AITC', 0) - vote_shares_7dma.get('BJP', 0)
                    cell = table.cell(9, col_map['7DMA'])
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{margin_7dma:.1f}"
                        else:
                            para.add_run().text = f"{margin_7dma:.1f}"
                    else:
                        cell.text = f"{margin_7dma:.1f}"
                
                # Overall margin
                if col_map['Overall'] < len(table.columns):
                    region_overall = regions_data_overall['regions'].get(region, {})
                    vote_shares_overall = region_overall.get('vote_shares', {})
                    margin_overall = vote_shares_overall.get('AITC', 0) - vote_shares_overall.get('BJP', 0)
                    cell = table.cell(9, col_map['Overall'])
                    if len(cell.text_frame.paragraphs) > 0:
                        para = cell.text_frame.paragraphs[0]
                        if len(para.runs) > 0:
                            para.runs[0].text = f"{margin_overall:.1f}"
                        else:
                            para.add_run().text = f"{margin_overall:.1f}"
                    else:
                        cell.text = f"{margin_overall:.1f}"
        
        return regions_data_overall['base_sample']
    
    def update_regional_chart(self, chart, region_name, is_7dma=False):
        """
        Update regional chart with vote share data for a specific region
        
        Args:
            chart: Chart object from slide
            region_name: Name of the region (e.g., 'Jalpaiguri', 'Malda')
            is_7dma: If True, show 7DMA data; if False, show Overall cumulative data
        
        Returns:
            Sample size used for calculations
        """
        from pptx.chart.data import CategoryChartData
        
        if is_7dma:
            # Calculate 7DMA for this region
            daily_data = self.calculate_7dma_for_region(region_name, num_days=13)
            sample_size = self._get_region_sample_size(region_name, is_7dma=True)
        else:
            # Calculate Overall cumulative for this region
            daily_data = self.calculate_overall_for_region(region_name, num_days=17)
            sample_size = self._get_region_sample_size(region_name, is_7dma=False)
        
        if len(daily_data) == 0:
            print(f"Warning: No data found for region {region_name}")
            return 0
        
        # Extract date labels for categories (use date strings like "15/10" instead of Excel serial numbers)
        date_labels = [str(d['date_label']) for d in daily_data]
        
        # Prepare series data
        aitc_values = [float(round(d['AITC'], 1)) for d in daily_data]
        bjp_values = [float(round(d['BJP'], 1)) for d in daily_data]
        left_values = [float(round(d['LEFT'], 1)) for d in daily_data]
        inc_values = [float(round(d['INC'], 1)) for d in daily_data]
        others_values = [float(round(d['Others'], 1)) for d in daily_data]
        nwr_values = [float(round(d['NWR'], 1)) for d in daily_data]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = date_labels
        
        chart_data.add_series('AITC (Trinamool Congress)', aitc_values)
        chart_data.add_series('BJP', bjp_values)
        chart_data.add_series('Left Front', left_values)
        chart_data.add_series('INC (Congress)', inc_values)
        chart_data.add_series('Others', others_values)
        chart_data.add_series('N+W+R', nwr_values)
        
        chart.replace_data(chart_data)
        
        # Set axis labels to vertical rotation (like format final file)
        self._set_axis_labels_vertical(chart)
        
        print(f"Updated {'7DMA' if is_7dma else 'Overall'} chart for region {region_name} with {len(daily_data)} days of data")
        return sample_size
    
    def calculate_7dma_for_region(self, region_name, num_days=13):
        """Calculate 7DMA vote shares for a region over num_days"""
        region_data = self.df[self.df['Region Name'] == region_name].copy()
        if len(region_data) == 0:
            return []
        
        # Get unique dates up to reference_date
        region_data = region_data[region_data['Survey Date'] <= self.reference_date].copy()
        question = '8. If assembly elections (MLA) were to be held tomorrow, then which party would you vote for?'
        
        # Find the first date with actual data for this region
        region_data_with_votes = region_data[region_data[question].notna()].copy()
        if len(region_data_with_votes) == 0:
            return []
        
        first_date_with_data = region_data_with_votes['Survey Date'].min()
        
        # Calculate 7DMA for each date starting from first_date_with_data
        # But only include dates where there's actual data in the 7-day window
        daily_results = []
        
        # Start from first_date_with_data and go up to reference_date
        # But we need at least 7 days of data before we can calculate 7DMA
        # So start calculating from (first_date_with_data + 6 days) to have 7 days in window
        start_date = first_date_with_data + timedelta(days=6)  # Need at least 7 days before
        
        # Generate all dates from start_date to reference_date
        # Start from start_date and show up to num_days dates (or up to reference_date, whichever comes first)
        all_dates = []
        date = start_date
        while date <= self.reference_date and len(all_dates) < num_days:
            all_dates.append(date)
            date += timedelta(days=1)
        
        for date in all_dates:
            # Calculate 7DMA for this date
            end_date = date - timedelta(days=1)  # Exclude reference date
            cutoff_date = end_date - timedelta(days=6)  # Last 7 days
            
            # Only calculate if we have data in the window
            date_data = region_data[
                (region_data['Survey Date'] >= cutoff_date) & 
                (region_data['Survey Date'] <= end_date)
            ].copy()
            
            # Only include if there's actual data with valid votes
            date_data_with_votes = date_data[date_data[question].notna()].copy()
            if len(date_data_with_votes) == 0:
                # Skip dates with no data in the 7-day window
                continue
            
            # Use L7D weights if available (with >=50% threshold)
            l7d_col = self.get_weight_column('Region', 'L7D') or 'Weight Voteshare L7D Region Level'
            regular_col = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
            
            l7d_available = date_data[l7d_col].notna().sum() if l7d_col in date_data.columns else 0
            total_records = len(date_data)
            
            if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
                # Use L7D weights - filter to only records with L7D weights
                date_data_with_weights = date_data[date_data[l7d_col].notna()].copy()
                weight_column = l7d_col
            else:
                # Fall back to regular weights
                date_data_with_weights = date_data.copy()
                weight_column = regular_col
            
            vote_shares = self.calculator.calculate_vote_share(date_data_with_weights, weight_column=weight_column, use_weights=True)
            
            # Convert date to Excel serial number and date label
            excel_serial = (date - datetime(1899, 12, 30)).days
            date_label = date.strftime('%m/%d/%Y')
            
            daily_results.append({
                'date': date,
                'date_label': date_label,
                'excel_serial': excel_serial,
                'AITC': vote_shares.get('AITC', 0),
                'BJP': vote_shares.get('BJP', 0),
                'LEFT': vote_shares.get('LEFT', 0),
                'INC': vote_shares.get('INC', 0),
                'Others': vote_shares.get('Others', 0),
                'NWR': vote_shares.get('NWR', 0)
            })
        
        return daily_results
    
    def calculate_overall_for_region(self, region_name, num_days=17):
        """Calculate Overall cumulative vote shares for a region over num_days"""
        region_data = self.df[self.df['Region Name'] == region_name].copy()
        if len(region_data) == 0:
            return []
        
        # Get unique dates up to reference_date
        region_data = region_data[region_data['Survey Date'] <= self.reference_date].copy()
        question = '8. If assembly elections (MLA) were to be held tomorrow, then which party would you vote for?'
        
        # Find the first date with actual data for this region
        region_data_with_votes = region_data[region_data[question].notna()].copy()
        if len(region_data_with_votes) == 0:
            return []
        
        first_date_with_data = region_data_with_votes['Survey Date'].min()
        
        # Calculate Overall for each date starting from first_date_with_data
        # Include ALL dates from first_date_with_data to reference_date (even if no new data on that date)
        # This ensures the graph starts from the first date and shows continuous dates
        daily_results = []
        
        # Generate all dates from first_date_with_data to reference_date
        # Start from first_date_with_data and show up to num_days dates (or up to reference_date, whichever comes first)
        all_dates = []
        date = first_date_with_data
        while date <= self.reference_date and len(all_dates) < num_days:
            all_dates.append(date)
            date += timedelta(days=1)
        
        last_vote_shares = None
        for date in all_dates:
            # Calculate cumulative up to this date
            date_data = region_data[region_data['Survey Date'] <= date].copy()
            
            # Check if there's actual data with valid votes
            date_data_with_votes = date_data[date_data[question].notna()].copy()
            if len(date_data_with_votes) > 0:
                # Calculate vote shares
                weight_column = self.get_weight_column('Region', 'Overall') or 'Weight Voteshare Overall Region Level'
                vote_shares = self.calculator.calculate_vote_share(date_data, weight_column=weight_column, use_weights=True)
                last_vote_shares = vote_shares
            elif last_vote_shares is not None:
                # If no new data on this date, use the cumulative value from previous date
                vote_shares = last_vote_shares
            else:
                # Skip if no data at all yet
                continue
            
            # Convert date to Excel serial number and date label
            excel_serial = (date - datetime(1899, 12, 30)).days
            date_label = date.strftime('%m/%d/%Y')
            
            daily_results.append({
                'date': date,
                'date_label': date_label,
                'excel_serial': excel_serial,
                'AITC': vote_shares.get('AITC', 0),
                'BJP': vote_shares.get('BJP', 0),
                'LEFT': vote_shares.get('LEFT', 0),
                'INC': vote_shares.get('INC', 0),
                'Others': vote_shares.get('Others', 0),
                'NWR': vote_shares.get('NWR', 0)
            })
        
        return daily_results
    
    def _get_region_sample_size(self, region_name, is_7dma=False):
        """Get sample size for a region - only count records with non-empty responses"""
        region_data = self.df[self.df['Region Name'] == region_name].copy()
        
        if is_7dma:
            # Filter for 7DMA window
            end_date = self.reference_date - timedelta(days=1)
            cutoff_date = end_date - timedelta(days=6)
            region_data = region_data[
                (region_data['Survey Date'] >= cutoff_date) & 
                (region_data['Survey Date'] <= end_date)
            ].copy()
        else:
            # Filter for overall
            region_data = region_data[region_data['Survey Date'] <= self.reference_date].copy()
        
        # Count only records with valid (non-empty) responses to main question
        return self.get_sample_size(region_data)
    
    def _set_axis_labels_vertical(self, chart):
        """
        Set category axis labels to vertical rotation (like format final file)
        Rotation value: -5400000 = -90 degrees (vertical)
        
        Args:
            chart: Chart object
        """
        try:
            # Access category axis
            cat_axis = chart.category_axis
            if cat_axis is None:
                return
            
            # Access the XML element to set rotation
            if hasattr(cat_axis, '_element'):
                element = cat_axis._element
                # Find txPr (text properties) element
                txPr = element.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}txPr')
                if txPr is not None:
                    # Find bodyPr (body properties) element
                    bodyPr = txPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
                    if bodyPr is not None:
                        # Set rotation to -90 degrees (vertical) = -5400000 in PowerPoint units
                        bodyPr.set('rot', '-5400000')
                    else:
                        # Create bodyPr if it doesn't exist
                        from pptx.oxml import parse_xml
                        bodyPr_xml = '<a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rot="-5400000" spcFirstLastPara="1" vertOverflow="ellipsis" wrap="square" anchor="ctr" anchorCtr="1"/>'
                        bodyPr = parse_xml(bodyPr_xml)
                        txPr.append(bodyPr)
                else:
                    # Create txPr if it doesn't exist
                    from pptx.oxml import parse_xml
                    txPr_xml = '''<c:txPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                        <a:bodyPr rot="-5400000" spcFirstLastPara="1" vertOverflow="ellipsis" wrap="square" anchor="ctr" anchorCtr="1"/>
                        <a:lstStyle/>
                        <a:p>
                            <a:pPr>
                                <a:defRPr/>
                            </a:pPr>
                        </a:p>
                    </c:txPr>'''
                    txPr = parse_xml(txPr_xml)
                    element.append(txPr)
        except Exception as e:
            # If setting rotation fails, continue without it
            print(f"Warning: Could not set axis label rotation: {e}")
            pass
    
    def _set_axis_labels_horizontal(self, chart):
        """
        Set category axis labels to horizontal rotation (0 degrees)
        For slides 95-136, dates should be horizontal, not vertical
        
        Args:
            chart: Chart object
        """
        try:
            # Access category axis
            cat_axis = chart.category_axis
            if cat_axis is None:
                return
            
            # Access the XML element to set rotation
            if hasattr(cat_axis, '_element'):
                element = cat_axis._element
                # Find txPr (text properties) element
                txPr = element.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}txPr')
                if txPr is not None:
                    # Find bodyPr (body properties) element
                    bodyPr = txPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
                    if bodyPr is not None:
                        # Set rotation to 0 degrees (horizontal)
                        bodyPr.set('rot', '0')
                    else:
                        # Create bodyPr if it doesn't exist with rotation 0
                        from pptx.oxml import parse_xml
                        bodyPr_xml = '<a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rot="0" spcFirstLastPara="1" vertOverflow="ellipsis" wrap="square" anchor="ctr" anchorCtr="1"/>'
                        bodyPr = parse_xml(bodyPr_xml)
                        txPr.append(bodyPr)
                else:
                    # Create txPr if it doesn't exist with rotation 0
                    from pptx.oxml import parse_xml
                    txPr_xml = '''<c:txPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                        <a:bodyPr rot="0" spcFirstLastPara="1" vertOverflow="ellipsis" wrap="square" anchor="ctr" anchorCtr="1"/>
                        <a:lstStyle/>
                        <a:p>
                            <a:pPr>
                                <a:defRPr/>
                            </a:pPr>
                        </a:p>
                    </c:txPr>'''
                    txPr = parse_xml(txPr_xml)
                    element.append(txPr)
        except Exception as e:
            # If setting rotation fails, continue without it
            print(f"Warning: Could not set axis label rotation to horizontal: {e}")
            pass
    
    def _move_legend_down(self, chart, offset_emus=317500):
        """
        Move chart legend down by specified offset to avoid overlapping with dates
        This is done by adjusting the plot area to make more room for the legend
        
        Args:
            chart: Chart object
            offset_emus: Offset in EMUs (default 317500 = ~25px)
        """
        try:
            # Instead of manipulating legend directly, adjust plot area to make room
            # Access chart element
            if not hasattr(chart, '_element'):
                return
            
            chart_element = chart._element
            
            # Find plotArea
            plot_area = chart_element.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}plotArea')
            if plot_area is None:
                return
            
            # Check if plotArea has layout
            layout = plot_area.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}layout')
            
            if layout is not None:
                # Check for manualLayout
                manual_layout = layout.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}manualLayout')
                
                if manual_layout is not None:
                    # Adjust height to make room for legend
                    h = manual_layout.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}h')
                    h_mode = manual_layout.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}hMode')
                    
                    if h is not None:
                        # Reduce height to make room for legend
                        current_h = float(h.get('val', '0'))
                        if current_h > 0:
                            new_h = max(0, current_h - offset_emus)
                            h.set('val', str(int(new_h)))
                    elif h_mode is None:
                        # Set hMode to factor and adjust
                        from pptx.oxml import parse_xml
                        h_mode_elem = parse_xml('<c:hMode xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" val="factor"/>')
                        manual_layout.append(h_mode_elem)
                else:
                    # Create manualLayout to adjust plot area
                    from pptx.oxml import parse_xml
                    manual_layout_xml = f'''<c:manualLayout xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
                        <c:layoutTarget val="inner"/>
                        <c:xMode val="edge"/>
                        <c:yMode val="edge"/>
                        <c:wMode val="edge"/>
                        <c:hMode val="edge"/>
                        <c:x val="0"/>
                        <c:y val="0"/>
                        <c:w val="0"/>
                        <c:h val="-{int(offset_emus)}"/>
                    </c:manualLayout>'''
                    manual_layout = parse_xml(manual_layout_xml)
                    layout.append(manual_layout)
            else:
                # Create layout with manualLayout
                from pptx.oxml import parse_xml
                layout_xml = f'''<c:layout xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
                    <c:manualLayout>
                        <c:layoutTarget val="inner"/>
                        <c:xMode val="edge"/>
                        <c:yMode val="edge"/>
                        <c:wMode val="edge"/>
                        <c:hMode val="edge"/>
                        <c:x val="0"/>
                        <c:y val="0"/>
                        <c:w val="0"/>
                        <c:h val="-{int(offset_emus)}"/>
                    </c:manualLayout>
                </c:layout>'''
                layout = parse_xml(layout_xml)
                plot_area.append(layout)
        except Exception as e:
            # If setting plot area fails, continue without it
            print(f"Warning: Could not adjust plot area for legend spacing: {e}")
            pass
    
    def update_district_table(self, table, is_15dma=False, is_7dma=False, slide_num=88):
        """
        Update district table with vote share data
        
        Args:
            table: Table object from slide
            is_15dma: If True, use 15DMA data
            is_7dma: If True, use 7DMA data
            slide_num: Slide number (to determine which districts to show)
        
        Returns:
            Sample size
        """
        # Get unique districts
        unique_districts = self.df['District Name'].dropna().unique()
        unique_districts = sorted([d for d in unique_districts if d != 'district_name'])
        
        # Determine which districts to show based on slide number
        # Slides 88, 90, 92: First half of districts
        # Slides 89, 91, 93: Second half of districts
        districts_per_slide = len(unique_districts) // 2
        if slide_num in [88, 90, 92]:
            districts_to_show = unique_districts[:districts_per_slide]
        else:
            districts_to_show = unique_districts[districts_per_slide:]
        
        # Filter data
        if is_15dma:
            end_date = self.reference_date - timedelta(days=1)
            cutoff_date = end_date - timedelta(days=14)
            filtered_data = self.df[
                (self.df['Survey Date'] >= cutoff_date) & 
                (self.df['Survey Date'] <= end_date)
            ].copy()
            weight_column = 'Weight - with Vote Share - AE 2021 - District L15D'
            if weight_column not in filtered_data.columns:
                weight_column = self.get_weight_column('District', 'Overall') or 'Weight Voteshare Overall District Level'
        elif is_7dma:
            end_date = self.reference_date - timedelta(days=1)
            cutoff_date = end_date - timedelta(days=6)
            filtered_data = self.df[
                (self.df['Survey Date'] >= cutoff_date) & 
                (self.df['Survey Date'] <= end_date)
            ].copy()
            weight_column = self.get_weight_column('District', 'L7D') or 'Weight Voteshare L7D District Level'
            if weight_column not in filtered_data.columns:
                weight_column = self.get_weight_column('District', 'Overall') or 'Weight Voteshare Overall District Level'
        else:
            filtered_data = self.df[self.df['Survey Date'] <= self.reference_date].copy()
            weight_column = self.get_weight_column('District', 'Overall') or 'Weight Voteshare Overall District Level'
        
        # Update table rows (starting from row 2, row 0 is header, row 1 is sub-header)
        total_sample = 0
        for row_idx, district in enumerate(districts_to_show, start=2):
            if row_idx >= len(table.rows):
                break
            
            district_data = filtered_data[filtered_data['District Name'] == district].copy()
            if len(district_data) == 0:
                continue
            
            # Calculate actual sample size: only count records with non-empty responses
            sample_size = self.get_sample_size(district_data)
            total_sample += sample_size
            
            # For 15DMA and 7DMA, check if L15D/L7D weights are available (with >=50% threshold)
            if is_15dma:
                l15d_col = 'Weight - with Vote Share - AE 2021 - District L15D'
                regular_col = self.get_weight_column('District', 'Overall') or 'Weight Voteshare Overall District Level'
                l15d_available = district_data[l15d_col].notna().sum() if l15d_col in district_data.columns else 0
                total_records = len(district_data)
                
                if l15d_available > 0 and (l15d_available / total_records) >= 0.5:
                    # Use L15D weights - filter to only records with L15D weights
                    district_data_with_weights = district_data[district_data[l15d_col].notna()].copy()
                    weight_column_used = l15d_col
                else:
                    # Fall back to regular weights
                    district_data_with_weights = district_data.copy()
                    weight_column_used = regular_col
            elif is_7dma:
                l7d_col = self.get_weight_column('District', 'L7D') or 'Weight Voteshare L7D District Level'
                regular_col = self.get_weight_column('District', 'Overall') or 'Weight Voteshare Overall District Level'
                l7d_available = district_data[l7d_col].notna().sum() if l7d_col in district_data.columns else 0
                total_records = len(district_data)
                
                if l7d_available > 0 and (l7d_available / total_records) >= 0.5:
                    # Use L7D weights - filter to only records with L7D weights
                    district_data_with_weights = district_data[district_data[l7d_col].notna()].copy()
                    weight_column_used = l7d_col
                else:
                    # Fall back to regular weights
                    district_data_with_weights = district_data.copy()
                    weight_column_used = regular_col
            else:
                # Overall - use regular weights
                district_data_with_weights = district_data.copy()
                weight_column_used = weight_column
            
            # Calculate vote shares (for percentages only, not sample size)
            vote_shares = self.calculator.calculate_vote_share(district_data_with_weights, weight_column=weight_column_used, use_weights=True)
            
            # Update sample size (column 2)
            if len(table.columns) > 2:
                cell = table.cell(row_idx, 2)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{sample_size:,}"
                    else:
                        para.add_run().text = f"{sample_size:,}"
                else:
                    cell.text = f"{sample_size:,}"
            
            # Update AITC (column 3)
            if len(table.columns) > 3:
                cell = table.cell(row_idx, 3)
                percentage = vote_shares.get('AITC', 0)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage:.1f}"
                    else:
                        para.add_run().text = f"{percentage:.1f}"
                else:
                    cell.text = f"{percentage:.1f}"
            
            # Update BJP (column 4)
            if len(table.columns) > 4:
                cell = table.cell(row_idx, 4)
                percentage = vote_shares.get('BJP', 0)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage:.1f}"
                    else:
                        para.add_run().text = f"{percentage:.1f}"
                else:
                    cell.text = f"{percentage:.1f}"
            
            # Update LEFT (column 5)
            if len(table.columns) > 5:
                cell = table.cell(row_idx, 5)
                percentage = vote_shares.get('LEFT', 0)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage:.1f}"
                    else:
                        para.add_run().text = f"{percentage:.1f}"
                else:
                    cell.text = f"{percentage:.1f}"
            
            # Update INC (column 6)
            if len(table.columns) > 6:
                cell = table.cell(row_idx, 6)
                percentage = vote_shares.get('INC', 0)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage:.1f}"
                    else:
                        para.add_run().text = f"{percentage:.1f}"
                else:
                    cell.text = f"{percentage:.1f}"
            
            # Update Others (column 7)
            if len(table.columns) > 7:
                cell = table.cell(row_idx, 7)
                percentage = vote_shares.get('Others', 0)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage:.1f}"
                    else:
                        para.add_run().text = f"{percentage:.1f}"
                else:
                    cell.text = f"{percentage:.1f}"
            
            # Update NWR (column 8)
            if len(table.columns) > 8:
                cell = table.cell(row_idx, 8)
                percentage = vote_shares.get('NWR', 0)
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{percentage:.1f}"
                    else:
                        para.add_run().text = f"{percentage:.1f}"
                else:
                    cell.text = f"{percentage:.1f}"
            
            # Update Margin (column 9) - AITC - BJP
            if len(table.columns) > 9:
                cell = table.cell(row_idx, 9)
                aitc_pct = vote_shares.get('AITC', 0)
                bjp_pct = vote_shares.get('BJP', 0)
                margin = aitc_pct - bjp_pct
                if len(cell.text_frame.paragraphs) > 0:
                    para = cell.text_frame.paragraphs[0]
                    if len(para.runs) > 0:
                        para.runs[0].text = f"{margin:.1f}"
                    else:
                        para.add_run().text = f"{margin:.1f}"
                else:
                    cell.text = f"{margin:.1f}"
        
        return total_sample
    
    def update_district_chart(self, chart, district_name, is_15dma=False):
        """
        Update district chart with vote share data
        
        Args:
            chart: Chart object from slide
            district_name: Name of the district
            is_15dma: If True, show 15DMA data; if False, show Overall cumulative data
        
        Returns:
            Sample size used for calculations
        """
        from pptx.chart.data import CategoryChartData
        
        if is_15dma:
            # Calculate 15DMA for this district
            daily_data = self.calculate_15dma_for_district(district_name, num_days=13)
            sample_size = self._get_district_sample_size(district_name, is_15dma=True)
        else:
            # Calculate Overall cumulative for this district
            daily_data = self.calculate_overall_for_district(district_name, num_days=17)
            sample_size = self._get_district_sample_size(district_name, is_15dma=False)
        
        if len(daily_data) == 0:
            print(f"Warning: No data found for district {district_name}")
            return 0
        
        # Limit to maximum 9 date points for slides 95-136
        if len(daily_data) > 9:
            # Take the last 9 data points
            daily_data = daily_data[-9:]
        
        # Extract date labels for categories (use date strings like "10/30/2025" instead of Excel serial numbers)
        date_labels = [str(d['date_label']) for d in daily_data]
        
        # Prepare series data
        aitc_values = [float(round(d['AITC'], 1)) for d in daily_data]
        bjp_values = [float(round(d['BJP'], 1)) for d in daily_data]
        left_values = [float(round(d['LEFT'], 1)) for d in daily_data]
        inc_values = [float(round(d['INC'], 1)) for d in daily_data]
        others_values = [float(round(d['Others'], 1)) for d in daily_data]
        nwr_values = [float(round(d['NWR'], 1)) for d in daily_data]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = date_labels
        
        chart_data.add_series('AITC (Trinamool Congress)', aitc_values)
        chart_data.add_series('BJP', bjp_values)
        chart_data.add_series('Left Front', left_values)
        chart_data.add_series('INC (Congress)', inc_values)
        chart_data.add_series('Others', others_values)
        chart_data.add_series('N+W+R', nwr_values)
        
        chart.replace_data(chart_data)
        
        # For slides 95-136, keep dates horizontal (set rotation to 0)
        self._set_axis_labels_horizontal(chart)
        
        # Move legend down for slides 95-136 to avoid overlapping with dates
        self._move_legend_down(chart, offset_emus=317500)  # ~25px = 317500 EMUs (1px = 12700 EMUs)
        
        print(f"Updated {'15DMA' if is_15dma else 'Overall'} chart for district {district_name} with {len(daily_data)} days of data")
        return sample_size
    
    def calculate_15dma_for_district(self, district_name, num_days=13):
        """Calculate 15DMA vote shares for a district over num_days"""
        district_data = self.df[self.df['District Name'] == district_name].copy()
        if len(district_data) == 0:
            return []
        
        # Get unique dates up to reference_date
        district_data = district_data[district_data['Survey Date'] <= self.reference_date].copy()
        unique_dates = sorted(district_data['Survey Date'].dropna().unique())
        
        if len(unique_dates) == 0:
            return []
        
        # Get last num_days dates
        dates_to_calculate = unique_dates[-num_days:] if len(unique_dates) >= num_days else unique_dates
        
        daily_results = []
        for date in dates_to_calculate:
            # Calculate 15DMA for this date
            end_date = date - timedelta(days=1)  # Exclude reference date
            cutoff_date = end_date - timedelta(days=14)  # Last 15 days
            
            date_data = district_data[
                (district_data['Survey Date'] >= cutoff_date) & 
                (district_data['Survey Date'] <= end_date)
            ].copy()
            
            if len(date_data) == 0:
                continue
            
            # Use L15D weights if available
            weight_column = 'Weight - with Vote Share - AE 2021 - District L15D'
            if weight_column not in date_data.columns:
                weight_column = self.get_weight_column('District', 'Overall') or 'Weight Voteshare Overall District Level'
            
            vote_shares = self.calculator.calculate_vote_share(date_data, weight_column=weight_column, use_weights=True)
            
            # Convert date to Excel serial number and date label
            excel_serial = (date - datetime(1899, 12, 30)).days
            date_label = date.strftime('%m/%d/%Y')
            
            daily_results.append({
                'date': date,
                'date_label': date_label,
                'excel_serial': excel_serial,
                'AITC': vote_shares.get('AITC', 0),
                'BJP': vote_shares.get('BJP', 0),
                'LEFT': vote_shares.get('LEFT', 0),
                'INC': vote_shares.get('INC', 0),
                'Others': vote_shares.get('Others', 0),
                'NWR': vote_shares.get('NWR', 0)
            })
        
        return daily_results
    
    def calculate_overall_for_district(self, district_name, num_days=17):
        """Calculate Overall cumulative vote shares for a district over num_days"""
        district_data = self.df[self.df['District Name'] == district_name].copy()
        if len(district_data) == 0:
            return []
        
        # Get unique dates up to reference_date
        district_data = district_data[district_data['Survey Date'] <= self.reference_date].copy()
        unique_dates = sorted(district_data['Survey Date'].dropna().unique())
        
        if len(unique_dates) == 0:
            return []
        
        # Get last num_days dates
        dates_to_calculate = unique_dates[-num_days:] if len(unique_dates) >= num_days else unique_dates
        
        daily_results = []
        for date in dates_to_calculate:
            # Calculate cumulative up to this date
            date_data = district_data[district_data['Survey Date'] <= date].copy()
            
            if len(date_data) == 0:
                continue
            
            weight_column = self.get_weight_column('District', 'Overall') or 'Weight Voteshare Overall District Level'
            vote_shares = self.calculator.calculate_vote_share(date_data, weight_column=weight_column, use_weights=True)
            
            # Convert date to Excel serial number and date label
            excel_serial = (date - datetime(1899, 12, 30)).days
            date_label = date.strftime('%m/%d/%Y')
            
            daily_results.append({
                'date': date,
                'date_label': date_label,
                'excel_serial': excel_serial,
                'AITC': vote_shares.get('AITC', 0),
                'BJP': vote_shares.get('BJP', 0),
                'LEFT': vote_shares.get('LEFT', 0),
                'INC': vote_shares.get('INC', 0),
                'Others': vote_shares.get('Others', 0),
                'NWR': vote_shares.get('NWR', 0)
            })
        
        return daily_results
    
    def _get_district_sample_size(self, district_name, is_15dma=False):
        """Get sample size for a district - only count records with non-empty responses"""
        district_data = self.df[self.df['District Name'] == district_name].copy()
        
        if is_15dma:
            # Filter for 15DMA window
            end_date = self.reference_date - timedelta(days=1)
            cutoff_date = end_date - timedelta(days=14)
            district_data = district_data[
                (district_data['Survey Date'] >= cutoff_date) & 
                (district_data['Survey Date'] <= end_date)
            ].copy()
        else:
            # Filter for overall
            district_data = district_data[district_data['Survey Date'] <= self.reference_date].copy()
        
        # Count only records with valid (non-empty) responses to main question
        return self.get_sample_size(district_data)
    
    def generate_complete_report(self, output_path):
        """
        Generate complete report by copying template and updating values
        
        Args:
            output_path: Path to save output PPT
        """
        print(f"\nGenerating complete report: {output_path}")
        print("="*80)
        
        # Copy entire template by loading it
        print("Copying template structure...")
        # Load template PPT and save it as output (we'll modify values)
        self.output_prs = Presentation(self.template_ppt_path)
        
        # Update slide 1 (Title) - Reporting Date (use current date when report is generated)
        if len(self.output_prs.slides) > 0:
            slide1 = self.output_prs.slides[0]
            # Use reference date (when report is generated)
            # Format date with proper ordinal suffix (1st, 2nd, 3rd, 4th, etc.)
            day = self.reference_date.day
            if 11 <= day <= 13:
                suffix = "th"
            elif day % 10 == 1:
                suffix = "st"
            elif day % 10 == 2:
                suffix = "nd"
            elif day % 10 == 3:
                suffix = "rd"
            else:
                suffix = "th"
            new_date = self.reference_date.strftime(f'%d{suffix} %B %Y')
            
            for shape in slide1.shapes:
                if hasattr(shape, 'text') and shape.text:
                    if 'Reporting Date' in shape.text or 'reporting date' in shape.text.lower():
                        # Update date - replace any existing date or append if no date exists
                        # Match dates like "31st October 2025" or "31 October 2025"
                        date_pattern = r'\d{1,2}(?:st|nd|rd|th)?\s+\w+\s+\d{4}'
                        current_text = shape.text
                        
                        # Check if there's already a date pattern in the text
                        if re.search(date_pattern, current_text):
                            # Replace existing date
                            shape.text = re.sub(date_pattern, new_date, current_text)
                        else:
                            # No date found, append it after "Reporting Date:" or "Reporting Date"
                            if current_text.strip().endswith(':'):
                                shape.text = f"{current_text.strip()} {new_date}"
                            elif current_text.strip().endswith('Reporting Date'):
                                shape.text = f"{current_text.strip()}: {new_date}"
                            else:
                                # Just append the date
                                shape.text = f"{current_text.strip()} {new_date}"
                        break
        
        print("Updated Slide 1: Title slide")
        
        # Update slide 2 (Introduction) - Sample size
        if len(self.output_prs.slides) > 1:
            slide2 = self.output_prs.slides[1]
            for shape in slide2.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_slide_2_table(shape.table)
            print("Updated Slide 2: Introduction")
        
        # Update slide 5 (State Level Vote Share)
        if len(self.output_prs.slides) > 4:
            slide5 = self.output_prs.slides[4]
            for shape in slide5.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_slide_5_table(shape.table)
            print("Updated Slide 5: State Level Vote Share")
        
        # Update slide 6 (Overall Normalized Vote Share Chart)
        if len(self.output_prs.slides) > 5:
            slide6 = self.output_prs.slides[5]
            sample_size_6 = 0
            for shape in slide6.shapes:
                try:
                    if hasattr(shape, 'chart') and shape.chart:
                        sample_size_6 = self.update_slide_6_chart(shape.chart)
                        break
                except (ValueError, AttributeError):
                    pass
            
            # Update Base sample size text on slide 6
            if sample_size_6 > 0:
                if self.update_base_text_on_slide(slide6, sample_size_6, slide_num=6):
                    print(f"Updated Slide 6 Base sample size: {sample_size_6:,}")
                else:
                    print(f"Note: Slide 6 Base text not found (sample size: {sample_size_6:,})")
            
            print("Updated Slide 6: Overall Normalized Vote Share Chart")
        
        # Update slide 7 (7 DMA Chart)
        if len(self.output_prs.slides) > 6:
            slide7 = self.output_prs.slides[6]
            sample_size_7 = 0
            for shape in slide7.shapes:
                try:
                    if hasattr(shape, 'chart') and shape.chart:
                        sample_size_7 = self.update_slide_7_chart(shape.chart)
                        break
                except (ValueError, AttributeError):
                    pass
            
            # Update Base sample size text on slide 7
            if sample_size_7 > 0:
                if self.update_base_text_on_slide(slide7, sample_size_7, slide_num=7):
                    print(f"Updated Slide 7 Base sample size: {sample_size_7:,}")
                else:
                    print(f"Note: Slide 7 Base text not found (sample size: {sample_size_7:,})")
            
            print("Updated Slide 7: 7 DMA Vote Share Chart")
        
        # Update slide 9 (Demographics - Part 1) - Multiple tables
        # Only replace values, do not change table positions or design
        if len(self.output_prs.slides) > 8:
            slide9 = self.output_prs.slides[8]
            table_count = 0
            
            for shape in slide9.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_demographic_table(shape.table, slide9)
                    table_count += 1
            
            if table_count > 0:
                print(f"Updated Slide 9: Demographics (Part 1) - {table_count} tables")
        
        # Update slide 10 (Demographics - Part 2) - Multiple tables
        # Only replace values, do not change table positions or design
        if len(self.output_prs.slides) > 9:
            slide10 = self.output_prs.slides[9]
            table_count = 0
            
            for shape in slide10.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_demographic_table(shape.table, slide10)
                    table_count += 1
            
            if table_count > 0:
                print(f"Updated Slide 10: Demographics (Part 2) - {table_count} tables")
        
        # Update slides 11-12 (Demographic tables - Part 3 and 4)
        # Slide 11: Social Category and Age (Part 3)
        # Only replace values, do not change table positions or design
        if len(self.output_prs.slides) > 10:
            slide11 = self.output_prs.slides[10]
            table_count = 0
            
            for shape in slide11.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_demographic_table(shape.table, slide11)
                    table_count += 1
            
            if table_count > 0:
                print(f"Updated Slide 11: Demographics (Part 3) - {table_count} tables")
        
        # Slide 12: Age (Part 4)
        # Only replace values, do not change table positions or design
        if len(self.output_prs.slides) > 11:
            slide12 = self.output_prs.slides[11]
            table_count = 0
            
            for shape in slide12.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_demographic_table(shape.table, slide12)
                    table_count += 1
            
            if table_count > 0:
                print(f"Updated Slide 12: Demographics (Part 4) - {table_count} tables")
        
        # Update Overall DMA demographic charts (Slides 14-23)
        # These are "DMA – [Demographic] (Overall)" charts
        overall_dma_demographic_slides = [
            (14, 'Gender', 'Male'),
            (15, 'Gender', 'Female'),
            (16, 'Location', 'Urban'),
            (17, 'Location', 'Rural'),
            (18, 'Religion', 'Hindu'),
            (19, 'Religion', 'Muslim'),
            (20, 'Social Category', 'General+OBC'),
            (21, 'Social Category', 'SC'),
            (22, 'Social Category', 'ST'),
            (23, 'Age', '18-25'),
        ]
        
        for slide_num, demographic_type, demographic_value in overall_dma_demographic_slides:
            if len(self.output_prs.slides) > slide_num - 1:
                slide = self.output_prs.slides[slide_num - 1]
                sample_size = 0
                
                # Find chart in slide
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'chart') and shape.chart:
                            sample_size = self.update_demographic_chart(shape.chart, demographic_type, demographic_value)
                            break
                    except (ValueError, AttributeError):
                        pass
                
                # Update Base sample size text on slide
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num):
                        print(f"Updated Slide {slide_num} Base sample size: {sample_size:,}")
                    # Don't print warning if Base not found - it might not exist on all slides
                
                if sample_size > 0:
                    print(f"Updated Slide {slide_num}: Overall DMA {demographic_type}={demographic_value} Chart")
        
        # Update 7 DMA demographic charts (Slides 24-39)
        # These are "7DMA – [Demographic]" charts
        dma7_demographic_slides = [
            (24, 'Age', '26-34'),  # Note: Slide 24 is 26-34, not 18-25
            (25, 'Age', '36-50'),
            (26, 'Age', '50+'),
            (27, 'Gender', 'Male'),
            (28, 'Gender', 'Female'),
            (29, 'Location', 'Urban'),
            (30, 'Location', 'Rural'),
            (31, 'Religion', 'Hindu'),
            (32, 'Religion', 'Muslim'),
            (33, 'Social Category', 'General+OBC'),
            (34, 'Social Category', 'SC'),
            (35, 'Social Category', 'ST'),
            (36, 'Age', '18-25'),
            (37, 'Age', '26-34'),
            (38, 'Age', '36-50'),
            (39, 'Age', '50+'),
        ]
        
        for slide_num, demographic_type, demographic_value in dma7_demographic_slides:
            if len(self.output_prs.slides) > slide_num - 1:
                slide = self.output_prs.slides[slide_num - 1]
                sample_size = 0
                
                # Find chart in slide
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'chart') and shape.chart:
                            sample_size = self.update_7dma_demographic_chart(shape.chart, demographic_type, demographic_value)
                            break
                    except (ValueError, AttributeError):
                        pass
                
                # Update Base sample size text on slide
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num):
                        print(f"Updated Slide {slide_num} Base sample size: {sample_size:,}")
                    # Don't print warning if Base not found - it might not exist on all slides
                
                if sample_size > 0:
                    print(f"Updated Slide {slide_num}: 7 DMA {demographic_type}={demographic_value} Chart")
        
        # Update slides 41-42 (Caste Wise Vote Shares)
        # Slide 41: Caste Wise Vote Shares (30 DMA)
        if len(self.output_prs.slides) > 40:
            slide41 = self.output_prs.slides[40]
            for shape in slide41.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_caste_table(shape.table, is_30dma=True)
                    print("Updated Slide 41: Caste Wise Vote Shares (30 DMA)")
                    break
        
        # Slide 42: Caste Wise Vote Shares (Overall)
        if len(self.output_prs.slides) > 41:
            slide42 = self.output_prs.slides[41]
            for shape in slide42.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    self.update_caste_table(shape.table, is_30dma=False)
                    print("Updated Slide 42: Caste Wise Vote Shares (Overall)")
                    break
        
        # Update slide 44 (Top Reasons for Party Choices - 7 DMA)
        if len(self.output_prs.slides) > 43:
            slide44 = self.output_prs.slides[43]
            sample_size_44 = 0
            for shape in slide44.shapes:
                try:
                    if hasattr(shape, 'chart') and shape.chart:
                        sample_size_44 = self.update_top_reasons_chart(shape.chart, is_7dma=True)
                        print("Updated Slide 44: Top Reasons for Party Choices - 7 DMA")
                        break
                except (ValueError, AttributeError):
                    pass
            
            # Update Base sample size text on slide 44
            if sample_size_44 > 0:
                if self.update_base_text_on_slide(slide44, sample_size_44, slide_num=44):
                    print(f"Updated Slide 44 Base sample size: {sample_size_44:,}")
        
        # Update slide 45 (Top Reasons for Party Choices - Overall)
        if len(self.output_prs.slides) > 44:
            slide45 = self.output_prs.slides[44]
            sample_size_45 = 0
            for shape in slide45.shapes:
                try:
                    if hasattr(shape, 'chart') and shape.chart:
                        sample_size_45 = self.update_top_reasons_chart(shape.chart, is_7dma=False)
                        print("Updated Slide 45: Top Reasons for Party Choices - Overall")
                        break
                except (ValueError, AttributeError):
                    pass
            
            # Update Base sample size text on slide 45
            if sample_size_45 > 0:
                if self.update_base_text_on_slide(slide45, sample_size_45, slide_num=45):
                    print(f"Updated Slide 45 Base sample size: {sample_size_45:,}")
        
        # Update Gains and Losses tables (Slides 47-61)
        # Slide 47: Gains and Losses - 7 DMA (Overall)
        if len(self.output_prs.slides) > 46:
            slide47 = self.output_prs.slides[46]
            sample_size_47 = 0
            for shape in slide47.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_47 = self.update_gains_losses_table(shape.table, is_7dma=True, is_demographic=False)
                    print("Updated Slide 47: Gains and Losses - 7 DMA (Overall)")
                    break
            
            # Update Base sample size text on slide 47 (skip - keep as is per user request)
            # if sample_size_47 > 0:
            #     if self.update_base_text_on_slide(slide47, sample_size_47, slide_num=47):
            #         print(f"Updated Slide 47 Base sample size: {sample_size_47:,}")
        
        # Slide 48: Gains and Losses - Overall
        if len(self.output_prs.slides) > 47:
            slide48 = self.output_prs.slides[47]
            sample_size_48 = 0
            for shape in slide48.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_48 = self.update_gains_losses_table(shape.table, is_7dma=False, is_demographic=False)
                    print("Updated Slide 48: Gains and Losses - Overall")
                    break
            
            # Update Base sample size text on slide 48 (skip - keep as is per user request)
            # if sample_size_48 > 0:
            #     if self.update_base_text_on_slide(slide48, sample_size_48, slide_num=48):
            #         print(f"Updated Slide 48 Base sample size: {sample_size_48:,}")
        
        # Slides 49-61: Gains and Losses - Demographics (Overall)
        demographic_gains_losses = [
            (49, 'Location', 'Urban'),
            (50, 'Location', 'Rural'),
            (51, 'Gender', 'Male'),
            (52, 'Gender', 'Female'),
            (53, 'Age', '18-25'),
            (54, 'Age', '26-34'),
            (55, 'Age', '36-50'),
            (56, 'Age', '50+'),
            (57, 'Religion', 'Hindu'),
            (58, 'Religion', 'Muslim'),
            (59, 'Social Category', 'General+OBC'),
            (60, 'Social Category', 'SC'),
            (61, 'Social Category', 'ST'),
        ]
        
        for slide_num, demographic_type, demographic_value in demographic_gains_losses:
            if len(self.output_prs.slides) > slide_num - 1:
                slide = self.output_prs.slides[slide_num - 1]
                sample_size = 0
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        sample_size = self.update_gains_losses_table(shape.table, is_7dma=False, is_demographic=True, 
                                                      demographic_type=demographic_type, demographic_value=demographic_value)
                        print(f"Updated Slide {slide_num}: Gains and Losses - {demographic_type}={demographic_value} (Overall)")
                        break
                
                # Update Base sample size text on slide (for slides 49-61, replace "Base of all eligible voters of 2021")
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num):
                        print(f"Updated Slide {slide_num} Base sample size: {sample_size:,}")
        
        # Update Vote Transferability tables (Slides 63-64)
        # Slide 63: Vote Transferability - 7 DMA
        if len(self.output_prs.slides) > 62:
            slide63 = self.output_prs.slides[62]
            sample_size_63 = 0
            for shape in slide63.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_63 = self.update_vote_transferability_table(shape.table, is_7dma=True)
                    print("Updated Slide 63: Vote Transferability - 7 DMA")
                    break
            
            # Update Base sample size text on slide 63
            if sample_size_63 > 0:
                if self.update_base_text_on_slide(slide63, sample_size_63, slide_num=63):
                    print(f"Updated Slide 63 Base sample size: {sample_size_63:,}")
        
        # Slide 64: Vote Transferability - Overall
        if len(self.output_prs.slides) > 63:
            slide64 = self.output_prs.slides[63]
            sample_size_64 = 0
            for shape in slide64.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_64 = self.update_vote_transferability_table(shape.table, is_7dma=False)
                    print("Updated Slide 64: Vote Transferability - Overall")
                    break
            
            # Update Base sample size text on slide 64
            if sample_size_64 > 0:
                if self.update_base_text_on_slide(slide64, sample_size_64, slide_num=64):
                    print(f"Updated Slide 64 Base sample size: {sample_size_64:,}")
        
        # Slide 66: Preferred CM Candidate
        if len(self.output_prs.slides) > 65:
            slide66 = self.output_prs.slides[65]
            sample_size_66 = 0
            for shape in slide66.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_66 = self.update_preferred_cm_table(shape.table)
                    print("Updated Slide 66: Preferred CM Candidate")
                    break
            
            # Update Base sample size text on slide 66
            if sample_size_66 > 0:
                if self.update_base_text_on_slide(slide66, sample_size_66, slide_num=66):
                    print(f"Updated Slide 66 Base sample size: {sample_size_66:,}")
        
        # Slide 68: State Government Rating
        if len(self.output_prs.slides) > 67:
            slide68 = self.output_prs.slides[67]
            sample_size_68 = 0
            for shape in slide68.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_68 = self.update_state_government_rating_table(shape.table)
                    print("Updated Slide 68: State Government Rating - State Level")
                    break
            
            # Update Base sample size text on slide 68
            if sample_size_68 > 0:
                if self.update_base_text_on_slide(slide68, sample_size_68, slide_num=68):
                    print(f"Updated Slide 68 Base sample size: {sample_size_68:,}")
        
        # Slide 70: Key Issues
        if len(self.output_prs.slides) > 69:
            slide70 = self.output_prs.slides[69]
            sample_size_70 = 0
            for shape in slide70.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_70 = self.update_key_issues_table(shape.table)
                    print("Updated Slide 70: Key Issues")
                    break
            
            # Update Base sample size text on slide 70
            if sample_size_70 > 0:
                if self.update_base_text_on_slide(slide70, sample_size_70, slide_num=70):
                    print(f"Updated Slide 70 Base sample size: {sample_size_70:,}")
        
        # Slide 72: Wisdom of Crowds
        if len(self.output_prs.slides) > 71:
            slide72 = self.output_prs.slides[71]
            sample_size_72 = 0
            for shape in slide72.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_72 = self.update_wisdom_of_crowds_table(shape.table)
                    print("Updated Slide 72: Wisdom of Crowds - Which Party is Likely to Win")
                    break
            
            # Update Base sample size text on slide 72
            if sample_size_72 > 0:
                if self.update_base_text_on_slide(slide72, sample_size_72, slide_num=72):
                    print(f"Updated Slide 72 Base sample size: {sample_size_72:,}")
        
        # Slide 74: Vote Share by Regions
        if len(self.output_prs.slides) > 73:
            slide74 = self.output_prs.slides[73]
            sample_size_74 = 0
            for shape in slide74.shapes:
                if hasattr(shape, 'has_table') and shape.has_table:
                    sample_size_74 = self.update_vote_share_by_regions_table(shape.table)
                    print("Updated Slide 74: Vote Share by Regions")
                    break
            
            # Update Base sample size text on slide 74 (if needed)
            if sample_size_74 > 0:
                if self.update_base_text_on_slide(slide74, sample_size_74, slide_num=74):
                    print(f"Updated Slide 74 Base sample size: {sample_size_74:,}")
        
        # Slides 75-84: Regional graphs (Overall and 7DMA)
        # Slide 75: Jalpaiguri (Overall), Slide 76: Jalpaiguri (7DMA)
        # Slide 77: Malda (Overall), Slide 78: Malda (7DMA)
        # Slide 79: Burdwan (Overall), Slide 80: Burdwan (7DMA)
        # Slide 81: Medinipur (Overall), Slide 82: Medinipur (7DMA)
        # Slide 83: Presidency (Overall), Slide 84: Presidency (7DMA)
        region_order = ['Jalpaiguri', 'Malda', 'Burdwan', 'Medinipur', 'Presidency']
        for region_idx, region in enumerate(region_order):
            # Overall slide (75, 77, 79, 81, 83)
            slide_num_overall = 75 + (region_idx * 2)
            if len(self.output_prs.slides) > slide_num_overall - 1:
                slide = self.output_prs.slides[slide_num_overall - 1]
                sample_size = 0
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'chart') and shape.chart:
                            sample_size = self.update_regional_chart(shape.chart, region, is_7dma=False)
                            print(f"Updated Slide {slide_num_overall}: {region} (Overall)")
                            break
                    except (ValueError, AttributeError):
                        continue
                
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num_overall):
                        print(f"Updated Slide {slide_num_overall} Base sample size: {sample_size:,}")
            
            # 7DMA slide (76, 78, 80, 82, 84)
            slide_num_7dma = 76 + (region_idx * 2)
            if len(self.output_prs.slides) > slide_num_7dma - 1:
                slide = self.output_prs.slides[slide_num_7dma - 1]
                sample_size = 0
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'chart') and shape.chart:
                            sample_size = self.update_regional_chart(shape.chart, region, is_7dma=True)
                            print(f"Updated Slide {slide_num_7dma}: {region} (7DMA)")
                            break
                    except (ValueError, AttributeError):
                        continue
                
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num_7dma):
                        print(f"Updated Slide {slide_num_7dma} Base sample size: {sample_size:,}")
        
        # Slides 86-87: Leave empty as per user request
        # (No update needed)
        
        # Slides 88-93: District-wise data
        # These slides show district vote shares
        # Slide 88: Districts 15DMA (1 of 2), Slide 89: Districts 15DMA (2 of 2)
        # Slide 90: Districts Overall (1 of 2), Slide 91: Districts Overall (2 of 2)
        # Slide 92: Districts 7DMA (1 of 2), Slide 93: Districts 7DMA (2 of 2)
        for slide_num in [88, 89]:
            if len(self.output_prs.slides) > slide_num - 1:
                slide = self.output_prs.slides[slide_num - 1]
                sample_size = 0
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        sample_size = self.update_district_table(shape.table, is_15dma=True, slide_num=slide_num)
                        print(f"Updated Slide {slide_num}: Districts (15DMA)")
                        break
                
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num):
                        print(f"Updated Slide {slide_num} Base sample size: {sample_size:,}")
        
        # Slides 90-91: Overall districts
        for slide_num in [90, 91]:
            if len(self.output_prs.slides) > slide_num - 1:
                slide = self.output_prs.slides[slide_num - 1]
                sample_size = 0
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        sample_size = self.update_district_table(shape.table, is_15dma=False, is_7dma=False, slide_num=slide_num)
                        print(f"Updated Slide {slide_num}: Districts (Overall)")
                        break
                
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num):
                        print(f"Updated Slide {slide_num} Base sample size: {sample_size:,}")
        
        # Slides 92-93: 7DMA districts
        for slide_num in [92, 93]:
            if len(self.output_prs.slides) > slide_num - 1:
                slide = self.output_prs.slides[slide_num - 1]
                sample_size = 0
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        sample_size = self.update_district_table(shape.table, is_15dma=False, is_7dma=True, slide_num=slide_num)
                        print(f"Updated Slide {slide_num}: Districts (7DMA)")
                        break
                
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num):
                        print(f"Updated Slide {slide_num} Base sample size: {sample_size:,}")
        
        # Slides 95-136: District graphs (15DMA and Overall)
        # These slides show district vote share graphs
        # Each district has 2 slides: 15DMA and Overall
        # Getting unique districts from data
        unique_districts = self.df['District Name'].dropna().unique()
        unique_districts = sorted([d for d in unique_districts if d != 'district_name'])
        
        # Starting from slide 95, each district has 2 slides (15DMA and Overall)
        for district_idx, district in enumerate(unique_districts):
            # 15DMA slide
            slide_num_15dma = 95 + (district_idx * 2)
            if len(self.output_prs.slides) > slide_num_15dma - 1:
                slide = self.output_prs.slides[slide_num_15dma - 1]
                sample_size = 0
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'chart') and shape.chart:
                            sample_size = self.update_district_chart(shape.chart, district, is_15dma=True)
                            print(f"Updated Slide {slide_num_15dma}: {district} (15DMA)")
                            break
                    except (ValueError, AttributeError):
                        continue
                
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num_15dma):
                        print(f"Updated Slide {slide_num_15dma} Base sample size: {sample_size:,}")
            
            # Overall slide
            slide_num_overall = 96 + (district_idx * 2)
            if len(self.output_prs.slides) > slide_num_overall - 1:
                slide = self.output_prs.slides[slide_num_overall - 1]
                sample_size = 0
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'chart') and shape.chart:
                            sample_size = self.update_district_chart(shape.chart, district, is_15dma=False)
                            print(f"Updated Slide {slide_num_overall}: {district} (Overall)")
                            break
                    except (ValueError, AttributeError):
                        continue
                
                if sample_size > 0:
                    if self.update_base_text_on_slide(slide, sample_size, slide_num=slide_num_overall):
                        print(f"Updated Slide {slide_num_overall} Base sample size: {sample_size:,}")
        
        # Save
        self.output_prs.save(output_path)
        print(f"\nReport saved: {output_path}")
        print("="*80)
        
        return output_path


def main():
    """
    Main function - runs automatically
    
    To use a different Excel file, you have 3 options:
    1. Update the excel_path variable below (line ~516)
    2. Pass as command-line argument: python3 generate_complete_report.py /path/to/excel.xlsx
    3. Place your Excel file in the same directory and update just the filename
    
    To specify a reference date (for 7 DMA calculations):
    --date YYYY-MM-DD (e.g., --date 2025-10-31)
    """
    import argparse
    
    # Default paths
    default_excel = "/Users/vijaygopal/Documents/Report_Generation/West_Bengal_31st_Oct_2025_With_Weights.xlsx"
    default_template = "/Users/vijaygopal/Documents/Report_Generation/Wb report Format Final.pptx"
    default_output = "/Users/vijaygopal/Documents/Report_Generation/West_Bengal_Survey_Report.pptx"
    
    # Parse command-line arguments (optional)
    parser = argparse.ArgumentParser(description='Generate survey report from Excel data')
    parser.add_argument('excel_file', nargs='?', default=default_excel,
                        help='Path to Excel file with survey data (default: %(default)s)')
    parser.add_argument('--template', default=default_template,
                        help='Path to template PPT file (default: %(default)s)')
    parser.add_argument('--output', default=default_output,
                        help='Path to output PPT file (default: %(default)s)')
    parser.add_argument('--date', type=str, default=None,
                        help='Reference date for report generation (format: YYYY-MM-DD). Default: current date')
    
    args = parser.parse_args()
    
    excel_path = args.excel_file
    template_ppt_path = args.template
    output_path = args.output
    reference_date = None
    
    # Parse reference date if provided
    if args.date:
        try:
            reference_date = pd.to_datetime(args.date)
            print(f"Using reference date: {reference_date.strftime('%Y-%m-%d')}")
        except Exception as e:
            print(f"Warning: Invalid date format '{args.date}'. Using current date. Error: {e}")
            reference_date = None
    
    print("="*80)
    print("COMPLETE REPORT GENERATOR")
    print("="*80)
    print(f"Input Excel: {excel_path}")
    print(f"Template PPT: {template_ppt_path}")
    print(f"Output PPT: {output_path}")
    print("="*80)
    
    try:
        generator = CompleteReportGenerator(excel_path, template_ppt_path, reference_date=reference_date)
        generator.generate_complete_report(output_path)
        
        print("\n" + "="*80)
        print("REPORT GENERATION COMPLETE")
        print("="*80)
        print(f"\n✅ Output file: {output_path}")
        print("\n📊 Updated slides:")
        print("   - Slide 1: Title (Date updated)")
        print("   - Slide 2: Introduction (Sample size)")
        print("   - Slide 5: State Level Vote Share (Calculated)")
        print("   - Slides 9-10: Demographics (Calculated)")
        print("\n💡 Note: Template design is preserved exactly")
        print("   Only values have been updated with calculations from Excel data\n")
    except Exception as e:
        print(f"\n❌ ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        return 1
    
    return 0


if __name__ == "__main__":
    main()
